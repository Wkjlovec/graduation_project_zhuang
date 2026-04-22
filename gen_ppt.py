# -*- coding: utf-8 -*-
"""Generate defense PPT with 8 slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2B, 0x7A, 0xE5)
LIGHT_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT = RGBColor(0xE8, 0x4D, 0x3D)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def add_para(text_frame, text, font_size=16, bold=False, color=TEXT_DARK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name=None):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    if font_name:
        p.font.name = font_name
    return p


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=14, font_color=TEXT_DARK, bold=False):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


# ============================================================
# SLIDE 1: Cover
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, BG_DARK)

add_textbox(s1, 1.5, 1.0, 10.3, 1.2,
            "\u57fa\u4e8e\u5fae\u670d\u52a1\u67b6\u6784\u7684\u9ad8\u6821\u8bba\u575b\u7cfb\u7edf\u8bbe\u8ba1\u4e0e\u5b9e\u73b0",
            font_size=36, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(s1, 1.5, 2.5, 10.3, 0.6,
            "Design and Implementation of a University Forum System Based on Microservice Architecture",
            font_size=16, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

info_lines = [
    "\u59d3\u540d\uff1a________    \u5b66\u53f7\uff1a________",
    "\u5bfc\u5e08\uff1a\u6c6a\u7ea2\u677e",
    "\u534e\u5357\u5e08\u8303\u5927\u5b66 \u00b7 \u8f6f\u4ef6\u5b66\u9662",
]
box = add_textbox(s1, 3.0, 3.8, 7.3, 2.0, "", font_size=18, color=TEXT_WHITE,
                  alignment=PP_ALIGN.CENTER)
box.text_frame.paragraphs[0].text = info_lines[0]
box.text_frame.paragraphs[0].font.size = Pt(20)
box.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
for line in info_lines[1:]:
    add_para(box.text_frame, line, font_size=20, color=TEXT_WHITE,
             alignment=PP_ALIGN.CENTER, space_before=Pt(12))

add_textbox(s1, 3.0, 6.2, 7.3, 0.5,
            "\u201c\u9488\u5bf9\u9ad8\u6821\u8bba\u575b\u5355\u4f53\u67b6\u6784\u6269\u5c55\u6027\u4e0d\u8db3\u3001"
            "\u7f3a\u4e4f\u79fb\u52a8\u7aef\u9002\u914d\u548c\u81ea\u52a8\u5316\u6d4b\u8bd5\u7684\u95ee\u9898\uff0c"
            "\u8bbe\u8ba1\u5e76\u5b9e\u73b0\u4e00\u5957\u57fa\u4e8e Spring Cloud \u7684\u5fae\u670d\u52a1\u8bba\u575b\u7cfb\u7edf\u201d",
            font_size=14, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Background & Innovation
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, BG_WHITE)

add_textbox(s2, 0.5, 0.3, 12.3, 0.7, "\u7814\u7a76\u80cc\u666f\u4e0e\u521b\u65b0\u70b9",
            font_size=30, bold=True, color=BG_DARK, alignment=PP_ALIGN.LEFT)

add_rect(s2, 0.5, 1.2, 5.8, 0.5, ACCENT,
         "\u73b0\u6709\u9ad8\u6821\u8bba\u575b\u7684\u4e09\u4e2a\u75db\u70b9",
         font_size=16, font_color=TEXT_WHITE, bold=True)

pain_box = add_textbox(s2, 0.7, 1.9, 5.5, 3.0, "", font_size=16, color=TEXT_DARK)
pains = [
    "\u2776 \u5355\u4f53\u67b6\u6784\u8026\u5408\uff1a\u5c40\u90e8\u4fee\u6539\u9700\u5168\u5c40\u91cd\u65b0\u90e8\u7f72\uff0c\u6269\u5c55\u56f0\u96be",
    "\u2777 \u7f3a\u4e4f\u79fb\u52a8\u7aef\uff1a\u4ec5\u6709 PC \u9875\u9762\uff0cCNNIC \u62a5\u544a\u663e\u793a\u624b\u673a\u4e0a\u7f51 99.6%",
    "\u2778 \u624b\u5de5\u6d4b\u8bd5\uff1a\u529f\u80fd\u9a8c\u8bc1\u4f9d\u8d56\u4eba\u5de5\u70b9\u51fb\uff0c\u7f3a\u5c11\u53ef\u91cf\u5316\u5224\u636e",
]
pain_box.text_frame.paragraphs[0].text = pains[0]
pain_box.text_frame.paragraphs[0].font.size = Pt(16)
for p in pains[1:]:
    add_para(pain_box.text_frame, p, font_size=16, space_before=Pt(16))

add_rect(s2, 7.0, 1.2, 5.8, 0.5, RGBColor(0x27, 0xAE, 0x60),
         "\u672c\u6587\u7684\u4e09\u4e2a\u521b\u65b0\u70b9",
         font_size=16, font_color=TEXT_WHITE, bold=True)

inno_box = add_textbox(s2, 7.2, 1.9, 5.5, 3.0, "", font_size=16, color=TEXT_DARK)
innos = [
    "\u2460 \u5fae\u670d\u52a1\u67b6\u6784\u5b8c\u6574\u843d\u5730\uff1aNacos + Gateway + OpenFeign",
    "\u2461 JWT + Redis \u6df7\u5408\u8ba4\u8bc1\uff1a\u4ee4\u724c\u4e3b\u52a8\u64a4\u9500 + \u524d\u7aef\u9759\u9ed8\u5237\u65b0",
    "\u2462 \u4e09\u5957\u81ea\u52a8\u5316\u6d4b\u8bd5\uff1a\u529f\u80fd\u56de\u5f52 / \u7f13\u5b58\u57fa\u51c6 / \u5e76\u53d1\u8d1f\u8f7d",
]
inno_box.text_frame.paragraphs[0].text = innos[0]
inno_box.text_frame.paragraphs[0].font.size = Pt(16)
for inn in innos[1:]:
    add_para(inno_box.text_frame, inn, font_size=16, space_before=Pt(16))


# ============================================================
# SLIDE 3: Requirements & Tech Route
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, BG_WHITE)

add_textbox(s3, 0.5, 0.3, 12.3, 0.7,
            "\u9700\u6c42\u5206\u6790\u4e0e\u6280\u672f\u8def\u7ebf",
            font_size=30, bold=True, color=BG_DARK)

add_rect(s3, 0.5, 1.2, 3.8, 0.5, ACCENT,
         "\u7528\u6237\u89d2\u8272", font_size=14, font_color=TEXT_WHITE, bold=True)
roles_box = add_textbox(s3, 0.7, 1.9, 3.5, 1.5, "", font_size=15)
roles_box.text_frame.paragraphs[0].text = "\u6e38\u5ba2\uff1a\u6d4f\u89c8\u3001\u641c\u7d22\u3001\u67e5\u770b\u516c\u544a"
add_para(roles_box.text_frame, "\u6ce8\u518c\u7528\u6237\uff1a\u53d1\u5e16\u3001\u8bc4\u8bba\u3001\u70b9\u8d5e\u3001\u901a\u77e5", font_size=15, space_before=Pt(8))
add_para(roles_box.text_frame, "\u6743\u9650\uff1a\u7f16\u8f91/\u5220\u9664\u4ec5\u9650\u4f5c\u8005\u672c\u4eba", font_size=15, space_before=Pt(8))

add_rect(s3, 0.5, 3.7, 3.8, 0.5, ACCENT,
         "\u516b\u5927\u529f\u80fd\u6a21\u5757", font_size=14, font_color=TEXT_WHITE, bold=True)
mods = add_textbox(s3, 0.7, 4.4, 3.5, 2.0, "", font_size=14)
mods.text_frame.paragraphs[0].text = "\u8ba4\u8bc1 \u00b7 \u5e16\u5b50 \u00b7 \u8bc4\u8bba \u00b7 \u70b9\u8d5e"
add_para(mods.text_frame, "\u901a\u77e5 \u00b7 \u641c\u7d22 \u00b7 \u5a92\u4f53 \u00b7 \u7f51\u5173", font_size=14, space_before=Pt(8))

add_rect(s3, 4.8, 1.2, 3.8, 0.5, RGBColor(0x27, 0xAE, 0x60),
         "\u975e\u529f\u80fd\u6027\u9700\u6c42", font_size=14, font_color=TEXT_WHITE, bold=True)
nfr = add_textbox(s3, 5.0, 1.9, 3.5, 3.0, "", font_size=14)
nfr_items = [
    "\u6027\u80fd\uff1a\u7f13\u5b58\u547d\u4e2d\u6beb\u79d2\u7ea7\uff0c\u672a\u547d\u4e2d <500ms",
    "\u5b89\u5168\uff1aJWT \u53cc\u4ee4\u724c + \u767b\u5f55\u5931\u8d25\u9501\u5b9a",
    "\u53ef\u6269\u5c55\uff1aNacos \u6ce8\u518c + \u7f51\u5173\u8def\u7531\u70ed\u63a5\u5165",
    "\u53ef\u7ef4\u62a4\uff1aDocker Compose + Bash \u542f\u505c",
    "\u517c\u5bb9\uff1aPC + \u79fb\u52a8\u7aef\u53cc\u7aef\u9002\u914d",
]
nfr.text_frame.paragraphs[0].text = nfr_items[0]
nfr.text_frame.paragraphs[0].font.size = Pt(14)
for item in nfr_items[1:]:
    add_para(nfr.text_frame, item, font_size=14, space_before=Pt(10))

add_rect(s3, 9.0, 1.2, 4.0, 0.5, RGBColor(0xF3, 0x9C, 0x12),
         "\u6280\u672f\u8def\u7ebf", font_size=14, font_color=TEXT_WHITE, bold=True)
route_items = [
    "\u2460 \u9700\u6c42\u5206\u6790\uff1a\u7528\u4f8b\u5efa\u6a21 + \u4e94\u7ef4\u6307\u6807",
    "\u2192",
    "\u2461 \u67b6\u6784\u8bbe\u8ba1\uff1a\u56db\u5c42\u62c6\u5206 + ER\u56fe",
    "\u2192",
    "\u2462 \u7f16\u7801\u5b9e\u73b0\uff1a\u5206\u670d\u52a1\u5f00\u53d1 + \u53cc\u7aef\u524d\u7aef",
    "\u2192",
    "\u2463 \u6d4b\u8bd5\u9a8c\u8bc1\uff1a\u529f\u80fd + \u7f13\u5b58 + \u5e76\u53d1",
]
route_box = add_textbox(s3, 9.2, 1.9, 3.7, 4.0, "", font_size=14)
route_box.text_frame.paragraphs[0].text = route_items[0]
route_box.text_frame.paragraphs[0].font.size = Pt(14)
for ri in route_items[1:]:
    sz = 20 if ri == "\u2192" else 14
    clr = ACCENT if ri == "\u2192" else TEXT_DARK
    add_para(route_box.text_frame, ri, font_size=sz, color=clr, space_before=Pt(6))


# ============================================================
# SLIDE 4: Architecture
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, BG_WHITE)

add_textbox(s4, 0.5, 0.3, 12.3, 0.7,
            "\u7cfb\u7edf\u67b6\u6784\uff1a\u56db\u5c42\u5fae\u670d\u52a1\u90e8\u7f72",
            font_size=30, bold=True, color=BG_DARK)

layers = [
    ("\u7528\u6237\u63a5\u5165\u5c42", "PC\u7aef Vue3+Element Plus :5173  |  \u79fb\u52a8\u7aef Vue3+Vant :5174",
     LIGHT_ACCENT, ACCENT),
    ("\u7f51\u5173\u5c42", "Spring Cloud Gateway :8080\u3000JWT\u9274\u6743 \u00b7 \u8def\u7531\u8f6c\u53d1 \u00b7 CORS",
     RGBColor(0xD5, 0xE8, 0xD4), RGBColor(0x27, 0xAE, 0x60)),
    ("\u4e1a\u52a1\u670d\u52a1\u5c42", "\u8ba4\u8bc1:8081  |  \u8bba\u575b:8082  |  \u901a\u77e5:8083  |  \u5a92\u4f53:8084  |  \u641c\u7d22:8085",
     RGBColor(0xFF, 0xF2, 0xCC), RGBColor(0xD6, 0xB6, 0x56)),
    ("\u57fa\u7840\u8bbe\u65bd\u5c42", "MySQL 8.4 :3306  |  Redis 7.4 :6379  |  Nacos 2.3.2 :8848",
     RGBColor(0xF0, 0xF0, 0xF0), TEXT_GRAY),
]

for i, (label, content, bg_clr, border_clr) in enumerate(layers):
    y = 1.3 + i * 1.4
    rect = add_rect(s4, 1.0, y, 11.3, 1.1, bg_clr, "", font_size=12)
    rect.line.color.rgb = border_clr
    rect.line.width = Pt(1.5)
    add_textbox(s4, 1.2, y + 0.05, 2.5, 0.4, label,
                font_size=14, bold=True, color=border_clr)
    add_textbox(s4, 1.2, y + 0.45, 10.8, 0.5, content,
                font_size=15, color=TEXT_DARK)

add_textbox(s4, 1.0, 6.5, 11.3, 0.5,
            "\u56fe 4-1 \u5fae\u670d\u52a1\u5206\u5c42\u90e8\u7f72\u4e0e\u57fa\u7840\u8bbe\u65bd\u7f16\u6392\u3000\u3000\u3000"
            "\u203b \u670d\u52a1\u95f4\u8c03\u7528\uff1a\u4ec5\u641c\u7d22\u670d\u52a1\u901a\u8fc7 OpenFeign \u8c03\u7528\u8bba\u575b\u670d\u52a1",
            font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: DB & API
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, BG_WHITE)

add_textbox(s5, 0.5, 0.3, 12.3, 0.7,
            "\u6570\u636e\u5e93\u4e0e\u63a5\u53e3\u8bbe\u8ba1",
            font_size=30, bold=True, color=BG_DARK)

add_rect(s5, 0.5, 1.2, 6.0, 0.5, ACCENT,
         "6 \u5f20\u4e1a\u52a1\u8868 \u00b7 5 \u4e2a\u5b9e\u4f53 \u00b7 7 \u6761\u5173\u7cfb",
         font_size=14, font_color=TEXT_WHITE, bold=True)

er_items = [
    "user_account \u2500\u2500 1:N \u2500\u2500 forum_post (creates)",
    "forum_section \u2500\u2500 1:N \u2500\u2500 forum_post (contains)",
    "forum_post \u2500\u2500 1:N \u2500\u2500 post_comment (has)",
    "forum_post \u2500\u2500 1:N \u2500\u2500 post_like (has)",
    "user_account \u2500\u2500 M:N \u2500\u2500 forum_post (\u70b9\u8d5e, via post_like)",
    "post_comment \u2500\u2500 1:N \u2500\u2500 post_comment (\u81ea\u5173\u8054\u56de\u590d)",
]
er_box = add_textbox(s5, 0.7, 1.9, 5.8, 4.0, "", font_size=13,
                     font_name="Consolas")
er_box.text_frame.paragraphs[0].text = er_items[0]
er_box.text_frame.paragraphs[0].font.size = Pt(13)
er_box.text_frame.paragraphs[0].font.name = "Consolas"
for item in er_items[1:]:
    p = add_para(er_box.text_frame, item, font_size=13, space_before=Pt(10))
    p.font.name = "Consolas"

add_rect(s5, 7.0, 1.2, 5.8, 0.5, RGBColor(0x27, 0xAE, 0x60),
         "22 \u4e2a RESTful API \u63a5\u53e3",
         font_size=14, font_color=TEXT_WHITE, bold=True)

api_items = [
    "\u8ba4\u8bc1\u670d\u52a1\uff1a6 \u4e2a\uff08\u6ce8\u518c/\u767b\u5f55/\u5237\u65b0/\u767b\u51fa/\u4e2a\u4eba\u4fe1\u606f/\u63a2\u6d4b\uff09",
    "\u8bba\u575b\u670d\u52a1\uff1a10 \u4e2a\uff08\u5e16\u5b50CRUD + \u8bc4\u8bba + \u70b9\u8d5e + \u5206\u533a\uff09",
    "\u901a\u77e5\u670d\u52a1\uff1a3 \u4e2a\uff08\u9996\u9875\u6458\u8981/\u5217\u8868/\u5df2\u8bfb\uff09",
    "\u641c\u7d22\u670d\u52a1\uff1a1 \u4e2a\uff08\u5173\u952e\u8bcd + \u5206\u533a + \u5206\u9875\uff09",
    "\u5a92\u4f53\u670d\u52a1\uff1a2 \u4e2a\uff08\u9996\u9875\u6458\u8981/\u63a8\u8350\u5217\u8868\uff09",
]
api_box = add_textbox(s5, 7.2, 1.9, 5.5, 3.5, "", font_size=14)
api_box.text_frame.paragraphs[0].text = api_items[0]
api_box.text_frame.paragraphs[0].font.size = Pt(14)
for item in api_items[1:]:
    add_para(api_box.text_frame, item, font_size=14, space_before=Pt(10))

add_textbox(s5, 7.2, 5.5, 5.5, 0.8,
            "\u6280\u672f\u9009\u578b\uff1aSpring Boot 3 + Spring Cloud Alibaba\n"
            "VS Django\uff1a\u7f3a\u5c11 Nacos/Gateway \u539f\u751f\u96c6\u6210",
            font_size=13, color=TEXT_GRAY)


# ============================================================
# SLIDE 6: Implementation
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, BG_WHITE)

add_textbox(s6, 0.5, 0.3, 12.3, 0.7,
            "\u6838\u5fc3\u5b9e\u73b0\uff1a\u8ba4\u8bc1\u95ed\u73af\u3001\u7f13\u5b58\u7b56\u7565\u4e0e\u5e76\u53d1\u63a7\u5236",
            font_size=28, bold=True, color=BG_DARK)

add_rect(s6, 0.5, 1.2, 4.0, 0.5, ACCENT,
         "JWT + Redis \u8ba4\u8bc1\u95ed\u73af",
         font_size=14, font_color=TEXT_WHITE, bold=True)
auth_code = (
    "login(username, password):\n"
    "  if locked(username): return ERROR\n"
    "  if !verify(password): \n"
    "    redis.incr(attempt); return ERROR\n"
    "  sid = UUID.random()\n"
    "  tokens = jwt.sign(uid, sid)\n"
    "  redis.set(sid -> uid, ttl=7d)\n"
    "  return tokens\n\n"
    "logout(token):\n"
    "  redis.del(session_id)\n"
    "  // \u4ee4\u724c\u7acb\u5373\u5931\u6548"
)
add_textbox(s6, 0.6, 1.9, 3.8, 3.5, auth_code,
            font_size=11, color=TEXT_DARK, font_name="Consolas")

add_rect(s6, 4.8, 1.2, 4.0, 0.5, RGBColor(0x27, 0xAE, 0x60),
         "Redis \u7f13\u5b58\u53cc\u8def\u5f84",
         font_size=14, font_color=TEXT_WHITE, bold=True)
cache_text = (
    "\u8bfb\u8def\u5f84\uff1a\n"
    "  \u8bf7\u6c42 \u2192 Redis \u547d\u4e2d? \u2192 \u8fd4\u56de\n"
    "       \u2514\u2500 \u672a\u547d\u4e2d \u2192 MySQL \u2192 \u5199\u5165Redis \u2192 \u8fd4\u56de\n\n"
    "\u5199\u8def\u5f84\uff1a\n"
    "  \u53d1\u5e03/\u7f16\u8f91/\u5220\u9664 \u2192 MySQL\u5199\u5165\n"
    "  \u2192 @CacheEvict \u6dd8\u6c70\u5217\u8868+\u8be6\u60c5\u7f13\u5b58\n\n"
    "\u8fc7\u671f\uff1a\u5217\u8868 2min / \u8be6\u60c5 5min / \u9ed8\u8ba4 10min"
)
add_textbox(s6, 4.9, 1.9, 3.8, 3.5, cache_text,
            font_size=12, color=TEXT_DARK)

add_rect(s6, 9.1, 1.2, 3.8, 0.5, HIGHLIGHT,
         "\u5e76\u53d1\u70b9\u8d5e\u53cc\u91cd\u53bb\u91cd",
         font_size=14, font_color=TEXT_WHITE, bold=True)
like_text = (
    "\u5e94\u7528\u5c42\uff1a\n"
    "  if exists(postId, userId):\n"
    "    return \"\u5df2\u70b9\u8d5e\"\n\n"
    "\u6570\u636e\u5e93\u5c42\uff1a\n"
    "  UNIQUE(post_id, user_id)\n"
    "  \u2192 \u5e76\u53d1\u63d2\u5165\u629b\u51fa\u5f02\u5e38\n"
    "  \u2192 catch \u8fd4\u56de\u63d0\u793a\n\n"
    "\u4e50\u89c2\u68c0\u67e5 + \u60b2\u89c2\u7ea6\u675f = \u96f6\u9501\u5f00\u9500"
)
add_textbox(s6, 9.2, 1.9, 3.6, 3.5, like_text,
            font_size=12, color=TEXT_DARK, font_name="Consolas")

add_textbox(s6, 0.5, 5.8, 12.3, 1.0,
            "\u5f00\u53d1\u5de5\u5177\u94fe\uff1a"
            "Java 17 + Maven 3.9 + Spring Boot 3.2.5 + Spring Cloud Alibaba 2023  |  "
            "Vue 3.5 + Vite 5.4 + TypeScript 5.9  |  "
            "Docker Compose + Bash + Python 3",
            font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Test Results
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, BG_WHITE)

add_textbox(s7, 0.5, 0.3, 12.3, 0.7,
            "\u6d4b\u8bd5\u7ed3\u679c",
            font_size=30, bold=True, color=BG_DARK)

add_rect(s7, 0.5, 1.2, 4.0, 0.5, ACCENT,
         "\u529f\u80fd\u6d4b\u8bd5", font_size=14, font_color=TEXT_WHITE, bold=True)
func_box = add_textbox(s7, 0.7, 1.9, 3.8, 2.5, "", font_size=15)
func_items = [
    "11 \u4e2a\u7aef\u5230\u7aef\u573a\u666f",
    "22 \u4e2a API \u63a5\u53e3\u5168\u8986\u76d6",
    "\u901a\u8fc7\u7387 100%",
    "\u542b 4 \u79cd\u5f02\u5e38\u8def\u5f84\u9a8c\u8bc1",
]
func_box.text_frame.paragraphs[0].text = func_items[0]
func_box.text_frame.paragraphs[0].font.size = Pt(15)
for fi in func_items[1:]:
    add_para(func_box.text_frame, fi, font_size=15, space_before=Pt(10))

add_rect(s7, 4.8, 1.2, 4.0, 0.5, RGBColor(0x27, 0xAE, 0x60),
         "\u7f13\u5b58\u57fa\u51c6", font_size=14, font_color=TEXT_WHITE, bold=True)
cache_box = add_textbox(s7, 5.0, 1.9, 3.8, 2.5, "", font_size=15)
cache_items = [
    "\u51b7\u542f\u52a8\uff1a59.12 ms",
    "\u70ed\u7f13\u5b58\uff1a14.98 ms (P95: 26.3ms)",
    "\u6027\u80fd\u63d0\u5347 74.66%",
    "\u70ed\u7f13\u5b58\u8bf7\u6c42 20 \u6b21",
]
cache_box.text_frame.paragraphs[0].text = cache_items[0]
cache_box.text_frame.paragraphs[0].font.size = Pt(15)
for ci in cache_items[1:]:
    add_para(cache_box.text_frame, ci, font_size=15, space_before=Pt(10))

add_rect(s7, 9.1, 1.2, 3.8, 0.5, RGBColor(0xF3, 0x9C, 0x12),
         "\u5e76\u53d1\u8d1f\u8f7d", font_size=14, font_color=TEXT_WHITE, bold=True)
load_box = add_textbox(s7, 9.3, 1.9, 3.6, 3.5, "", font_size=13)
load_items = [
    "\u5e76\u53d1   QPS      P95     \u6210\u529f\u7387",
    "  10    1186    10ms    100%",
    "  50    1420    51ms    100%",
    " 100    1392    97ms    99.95%",
    " 500    1557   226ms    100%",
    "1000    1450   366ms    99.99%",
]
load_box.text_frame.paragraphs[0].text = load_items[0]
load_box.text_frame.paragraphs[0].font.size = Pt(13)
load_box.text_frame.paragraphs[0].font.name = "Consolas"
load_box.text_frame.paragraphs[0].font.bold = True
for li in load_items[1:]:
    p = add_para(load_box.text_frame, li, font_size=13, space_before=Pt(6))
    p.font.name = "Consolas"


# ============================================================
# SLIDE 8: Conclusion
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s8, BG_DARK)

add_textbox(s8, 0.5, 0.5, 12.3, 0.7,
            "\u603b\u7ed3\u4e0e\u5c55\u671b",
            font_size=30, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_rect(s8, 1.0, 1.5, 5.5, 0.5, ACCENT,
         "\u4e09\u70b9\u6838\u5fc3\u7ed3\u8bba",
         font_size=16, font_color=TEXT_WHITE, bold=True)
conc_box = add_textbox(s8, 1.2, 2.2, 5.2, 2.5, "", font_size=15, color=TEXT_WHITE)
concs = [
    "\u2714 \u5fae\u670d\u52a1\u67b6\u6784\u5728\u9ad8\u6821\u8bba\u575b\u573a\u666f\u53ef\u5b9e\u9645\u843d\u5730",
    "\u2714 PC + \u79fb\u52a8\u7aef\u53cc\u7aef\u9002\u914d\u65b9\u6848\u53ef\u884c",
    "\u2714 \u81ea\u52a8\u5316\u6d4b\u8bd5\u53ef\u590d\u73b0\u3001\u7ed3\u679c\u53ef\u91cf\u5316",
]
conc_box.text_frame.paragraphs[0].text = concs[0]
conc_box.text_frame.paragraphs[0].font.size = Pt(15)
conc_box.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
for c in concs[1:]:
    add_para(conc_box.text_frame, c, font_size=15, color=TEXT_WHITE, space_before=Pt(14))

add_rect(s8, 7.0, 1.5, 5.5, 0.5, RGBColor(0xE8, 0x4D, 0x3D),
         "\u4e0d\u8db3\u4e0e\u5c55\u671b",
         font_size=16, font_color=TEXT_WHITE, bold=True)
limit_box = add_textbox(s8, 7.2, 2.2, 5.2, 2.5, "", font_size=15, color=TEXT_WHITE)
limits = [
    "\u25b3 \u641c\u7d22\u670d\u52a1\u91c7\u7528\u5185\u5b58\u8fc7\u6ee4\uff0c\u540e\u7eed\u5f15\u5165 Elasticsearch",
    "\u25b3 \u672a\u7ecf\u4e91\u7aef\u751f\u4ea7\u9a8c\u8bc1\uff0c\u540e\u7eed\u90e8\u7f72\u81f3 K8s \u96c6\u7fa4",
]
limit_box.text_frame.paragraphs[0].text = limits[0]
limit_box.text_frame.paragraphs[0].font.size = Pt(15)
limit_box.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
for lm in limits[1:]:
    add_para(limit_box.text_frame, lm, font_size=15, color=TEXT_WHITE, space_before=Pt(14))

add_textbox(s8, 1.0, 5.5, 11.3, 1.2,
            "\u4ee5\u4e0a\u662f\u6211\u7684\u6c47\u62a5\uff0c\u8bf7\u5404\u4f4d\u8001\u5e08\u6279\u8bc4\u6307\u6b63\u3002",
            font_size=24, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(s8, 1.0, 6.5, 11.3, 0.5,
            "Thank you!",
            font_size=20, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

# Save
prs.save("docs/\u7b54\u8fa9PPT.pptx")
print("PPT saved: docs/\u7b54\u8fa9PPT.pptx")
print(f"Total slides: {len(prs.slides)}")
