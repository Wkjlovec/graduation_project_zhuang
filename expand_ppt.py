# -*- coding: utf-8 -*-
"""Expand PPT slides 3/4/5 with additional content pages."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

BG_DARK = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
ACCENT = RGBColor(0x2B, 0x7A, 0xE5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE8, 0x4D, 0x3D)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=TEXT_DARK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    if font_name:
        p.font.name = font_name
    return p


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=14, font_color=TEXT_DARK, bold=False):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def make_titled_content(prs, title, items, accent_color=ACCENT, insert_idx=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_rect(slide, 0, 0, 13.333, 0.9, accent_color,
             title, font_size=22, font_color=TEXT_WHITE, bold=True)
    for i, (heading, desc) in enumerate(items):
        y = 1.2 + i * 1.05
        add_textbox(slide, 1.0, y, 11.3, 0.35, heading,
                    font_size=18, bold=True, color=accent_color)
        add_textbox(slide, 1.0, y + 0.38, 11.3, 0.55, desc,
                    font_size=14, color=TEXT_GRAY)
    if insert_idx is not None:
        slides = list(prs.slides._sldIdLst)
        new_slide_elem = slides[-1]
        prs.slides._sldIdLst.remove(new_slide_elem)
        prs.slides._sldIdLst.insert(insert_idx, new_slide_elem)
    return slide


prs = Presentation("docs/\u7b54\u8fa9PPT.pptx")

# ============================================================
# Chapter 3 expansion: insert after Slide 5 (index 4)
# New slide: Functional requirements detail
# ============================================================
make_titled_content(prs,
    "\u529f\u80fd\u6027\u9700\u6c42\uff1a\u516b\u4e2a\u6a21\u5757\u7684\u6838\u5fc3\u7528\u4f8b",
    [
        ("\u7528\u6237\u8ba4\u8bc1", "\u6ce8\u518c\u3001\u767b\u5f55\u3001\u4ee4\u724c\u5237\u65b0\u3001\u767b\u51fa\u56db\u4e2a\u63a5\u53e3\u5f62\u6210\u95ed\u73af\uff0c\u767b\u5f55\u5931\u8d25 5 \u6b21\u9501\u5b9a 10 \u5206\u949f"),
        ("\u5e16\u5b50\u7ba1\u7406", "\u53d1\u5e03\u3001\u6d4f\u89c8\u3001\u7f16\u8f91\u3001\u5220\u9664\uff0c\u5217\u8868\u652f\u6301\u5206\u9875\u548c\u5206\u533a\u7b5b\u9009\uff0c\u7f16\u8f91/\u5220\u9664\u4ec5\u9650\u4f5c\u8005"),
        ("\u8bc4\u8bba\u4e92\u52a8", "\u652f\u6301\u4e8c\u7ea7\u56de\u590d\u7ed3\u6784\uff08parent_comment_id\uff09\uff0c\u5e16\u5b50\u4f5c\u8005\u53ef\u5220\u9664\u81ea\u5df1\u5e16\u5b50\u4e0b\u7684\u4efb\u610f\u8bc4\u8bba"),
        ("\u70b9\u8d5e", "\u540c\u4e00\u7528\u6237\u5bf9\u540c\u4e00\u5e16\u5b50\u6700\u591a\u4e00\u6b21\u6709\u6548\u70b9\u8d5e\uff0c\u8054\u5408\u552f\u4e00\u7ea6\u675f + \u5e94\u7528\u5c42\u53cc\u91cd\u53bb\u91cd"),
        ("\u901a\u77e5 / \u641c\u7d22 / \u5a92\u4f53", "\u901a\u77e5\u652f\u6301\u5df2\u8bfb\u6807\u8bb0\uff0c\u641c\u7d22\u8de8\u670d\u52a1\u8c03\u7528\u8bba\u575b\u83b7\u53d6\u6570\u636e\uff0c\u5a92\u4f53\u9884\u7f6e\u97f3\u4e50/\u4e66\u7c4d\u63a8\u8350"),
    ],
    ACCENT, insert_idx=5)

# New slide: Non-functional requirements detail
make_titled_content(prs,
    "\u975e\u529f\u80fd\u6027\u9700\u6c42\uff1a\u4e94\u4e2a\u7ef4\u5ea6\u7684\u91cf\u5316\u6307\u6807",
    [
        ("\u6027\u80fd", "\u7f13\u5b58\u547d\u4e2d\u6beb\u79d2\u7ea7\u54cd\u5e94\uff0c\u672a\u547d\u4e2d <500ms\uff1b\u5e16\u5b50\u5217\u8868\u7f13\u5b58 2min\u3001\u8be6\u60c5 5min\u3001\u9ed8\u8ba4 10min"),
        ("\u5b89\u5168", "JWT \u53cc\u4ee4\u724c\uff08access 24h + refresh 7d\uff09\uff0cHMAC-SHA \u7b7e\u540d\uff0c\u767b\u5f55\u5931\u8d25 5 \u6b21\u9501\u5b9a 600s\uff0c\u7f51\u5173\u53cc\u91cd\u6821\u9a8c"),
        ("\u53ef\u6269\u5c55", "\u65b0\u6a21\u5757\u6ce8\u518c Nacos + \u7f51\u5173\u52a0\u4e00\u6761\u8def\u7531\u5373\u53ef\u63a5\u5165\uff0c\u65e0\u9700\u6539\u52a8\u5df2\u6709\u670d\u52a1\u4ee3\u7801"),
        ("\u53ef\u7ef4\u62a4", "Docker Compose \u4e00\u952e\u62c9\u8d77\u57fa\u7840\u8bbe\u65bd\uff0cBash \u811a\u672c\u6309\u987a\u5e8f\u542f\u505c + \u5065\u5eb7\u63a2\u6d4b"),
        ("\u517c\u5bb9", "PC \u7aef Element Plus + \u79fb\u52a8\u7aef Vant\uff0c\u4e24\u7aef\u5171\u4eab API \u5c42\u548c\u72b6\u6001\u7ba1\u7406"),
    ],
    GREEN, insert_idx=6)

# ============================================================
# Chapter 4 expansion: insert after current architecture slides
# New slide: Security design detail (after Slide 8, now shifted)
# ============================================================
# After expansion ch3 adds 2 slides, so old slide 8 (DB) is now at index 9
# Insert security slide at index 10
make_titled_content(prs,
    "\u5b89\u5168\u8bbe\u8ba1\uff1aJWT + Redis + \u7f51\u5173\u4e09\u9053\u9632\u7ebf",
    [
        ("\u7b2c\u4e00\u9053\uff1a\u5bc6\u7801\u52a0\u5bc6", "BCrypt \u54c8\u5e0c\u5b58\u50a8\uff0c\u539f\u59cb\u5bc6\u7801\u4e0d\u5728\u4efb\u4f55\u73af\u8282\u4ee5\u660e\u6587\u4fdd\u7559"),
        ("\u7b2c\u4e8c\u9053\uff1a\u4ee4\u724c\u7b7e\u53d1\u4e0e\u4f1a\u8bdd", "\u767b\u5f55\u6210\u529f\u540e\u7b7e\u53d1 access + refresh \u53cc\u4ee4\u724c\uff0cRedis \u5b58\u50a8 sessionId\u2192userId \u6620\u5c04"),
        ("\u7b2c\u4e09\u9053\uff1a\u7f51\u5173\u53cc\u91cd\u6821\u9a8c", "JwtRelayFilter \u5148\u9a8c JWT \u7b7e\u540d\u518d\u67e5 Redis \u4f1a\u8bdd\uff0c\u901a\u8fc7\u540e\u900f\u4f20 userId/username \u5230\u4e0b\u6e38"),
        ("\u767b\u5f55\u5931\u8d25\u9632\u62a4", "Redis \u539f\u5b50\u8ba1\u6570\uff0c\u8fde\u7eed 5 \u6b21\u5931\u8d25\u9501\u5b9a 600 \u79d2\uff0c\u9632\u66b4\u529b\u7834\u89e3"),
        ("\u767b\u51fa\u5373\u65f6\u5931\u6548", "\u5220\u9664 Redis \u4f1a\u8bdd\u8bb0\u5f55\uff0c\u5df2\u7b7e\u53d1\u7684\u4ee4\u724c\u5373\u4f7f\u672a\u8fc7\u671f\u4e5f\u4f1a\u5728\u7f51\u5173\u5c42\u88ab\u62d2\u7edd"),
    ],
    RED, insert_idx=11)

# New slide: Frontend dual-end implementation
make_titled_content(prs,
    "\u524d\u7aef\u53cc\u7aef\u5b9e\u73b0\uff1a\u5171\u4eab\u5e95\u5c42\u3001\u5404\u81ea\u9002\u914d",
    [
        ("PC \u7aef\uff08Element Plus\uff09", "\u5de6\u53f3\u5206\u680f\u5e03\u5c40\uff0c\u5de6\u4fa7\u5e16\u5b50\u5217\u8868 + \u641c\u7d22 + \u5206\u533a\u7b5b\u9009\uff0c\u53f3\u4fa7\u901a\u77e5\u6458\u8981 + \u5a92\u4f53\u63a8\u8350"),
        ("\u79fb\u52a8\u7aef\uff08Vant\uff09", "\u7eb5\u5411\u5806\u53e0\u5e03\u5c40\u9002\u914d\u5c0f\u5c4f\uff0c\u5e95\u90e8\u56fa\u5b9a\u5bfc\u822a\uff0c\u89e6\u5c4f\u4ea4\u4e92\u4f18\u5316"),
        ("\u5171\u4eab\u5c42", "Axios \u8bf7\u6c42\u5c42 + Pinia \u72b6\u6001\u7ba1\u7406 + Vue Router \u8def\u7531\u5b88\u536b\uff0c\u4fee\u6539\u4e00\u5904\u63a5\u53e3\u4e24\u7aef\u540c\u6b65\u751f\u6548"),
        ("\u9759\u9ed8\u5237\u65b0\u673a\u5236", "\u54cd\u5e94\u62e6\u622a\u5668\u5904\u7406 401\uff0c\u5e76\u53d1\u8bf7\u6c42\u52a0\u5165\u7b49\u5f85\u961f\u5217\uff0c\u5237\u65b0\u6210\u529f\u540e\u7edf\u4e00\u91ca\u653e\u91cd\u8bd5"),
        ("\u8def\u7531\u5b88\u536b", "\u53d1\u5e16\u9875\u3001\u901a\u77e5\u9875\u3001\u4e2a\u4eba\u4e2d\u5fc3\u8bbe\u7f6e\u5bfc\u822a\u5b88\u536b\uff0c\u672a\u767b\u5f55\u81ea\u52a8\u91cd\u5b9a\u5411\u5230\u767b\u5f55\u9875"),
    ],
    ACCENT, insert_idx=13)

# ============================================================
# Chapter 5 expansion: insert after test overview
# New slide: Functional test detail
# ============================================================
# Test overview is now further shifted, find it
# After adding 4 slides above, old slide 12 is now at ~index 16
make_titled_content(prs,
    "\u529f\u80fd\u6d4b\u8bd5\uff1a11 \u4e2a\u573a\u666f\u7684\u8986\u76d6\u8303\u56f4",
    [
        ("\u6b63\u5e38\u4e1a\u52a1\u6d41\u7a0b\uff08\u573a\u666f 1\u20135\uff09", "\u5206\u533a\u67e5\u8be2 \u2192 \u6ce8\u518c\u767b\u5f55 \u2192 \u53d1\u5e16\u8bc4\u8bba\u641c\u7d22 \u2192 \u901a\u77e5\u5a92\u4f53\u9996\u9875 \u2192 \u4e2a\u4eba\u901a\u77e5\u5df2\u8bfb"),
        ("\u5f02\u5e38\u8def\u5f84\u9a8c\u8bc1\uff08\u573a\u666f 6\u20139\uff09", "\u65e0\u4ee4\u724c\u53d1\u5e16\u21924010 | \u8d8a\u6743\u5220\u5e16\u21924031 | \u5e76\u53d1\u91cd\u590d\u70b9\u8d5e\u21920+4091 | \u767b\u51fa\u540e\u8bbf\u95ee\u21924010"),
        ("\u4ee4\u724c\u751f\u547d\u5468\u671f\uff08\u573a\u666f 10\u201311\uff09", "\u5237\u65b0\u4ee4\u724c\u83b7\u53d6\u65b0\u4ee4\u724c\u5bf9 | \u4ee4\u724c\u8fc7\u671f\u540e\u8bbf\u95ee\u53d7\u4fdd\u62a4\u63a5\u53e3\u88ab\u62e6\u622a"),
        ("API \u8986\u76d6", "\u5168\u90e8 22 \u4e2a\u63a5\u53e3\u81f3\u5c11\u88ab\u8c03\u7528\u4e00\u6b21\uff1a\u8ba4\u8bc1 6 + \u8bba\u575b 10 + \u901a\u77e5 3 + \u641c\u7d22 1 + \u5a92\u4f53 2"),
        ("\u6d4b\u8bd5\u65b9\u6cd5", "Python \u811a\u672c\u9a71\u52a8\u7684\u9ed1\u76d2\u7aef\u5230\u7aef\u6d4b\u8bd5\uff0c\u6bd4\u5bf9 HTTP \u72b6\u6001\u7801\u548c\u4e1a\u52a1\u7801\uff0c\u81ea\u52a8\u751f\u6210 Markdown \u62a5\u544a"),
    ],
    ACCENT, insert_idx=17)

# New slide: Performance test detail
make_titled_content(prs,
    "\u6027\u80fd\u6d4b\u8bd5\uff1a\u7f13\u5b58\u6548\u679c\u4e0e\u5e76\u53d1\u627f\u8f7d\u80fd\u529b",
    [
        ("\u7f13\u5b58\u57fa\u51c6\u6d4b\u8bd5\u539f\u7406", "\u5148\u53d1\u4e00\u6b21\u51b7\u542f\u52a8\u8bf7\u6c42\uff08\u7a7f\u900f\u5230 MySQL\uff09\uff0c\u518d\u53d1 20 \u6b21\u70ed\u7f13\u5b58\u8bf7\u6c42\uff08\u547d\u4e2d Redis\uff09\uff0c\u5bf9\u6bd4\u54cd\u5e94\u65f6\u95f4"),
        ("\u7f13\u5b58\u6548\u679c", "\u51b7\u542f\u52a8 59.12ms \u2192 \u70ed\u7f13\u5b58 14.98ms\uff0cP95 \u4e3a 26.3ms\uff0c\u6027\u80fd\u63d0\u5347 74.66%"),
        ("\u5e76\u53d1\u8d1f\u8f7d\u6d4b\u8bd5\u539f\u7406", "ThreadPoolExecutor \u521b\u5efa\u7ebf\u7a0b\u6c60\uff0c\u6bcf\u4e2a\u865a\u62df\u7528\u6237\u53d1 20 \u6b21 GET \u8bf7\u6c42\uff0c\u5206\u4e94\u6863\u9010\u6b65\u52a0\u538b"),
        ("\u5e76\u53d1\u7ed3\u679c", "10\u219d1186 QPS | 100\u219d1392 QPS, 99.95% | 500\u219d1557 QPS | 1000\u219d1450 QPS, 99.99%"),
        ("\u5c40\u9650\u6027\u8bf4\u660e", "\u672c\u5730\u5355\u673a\u6d4b\u8bd5\uff0c\u7edd\u5bf9\u6570\u503c\u4e0e\u751f\u4ea7\u73af\u5883\u6709\u5dee\u5f02\uff0c\u4f46\u7f13\u5b58\u63d0\u5347\u6bd4\u4f8b\u548c\u5e76\u53d1\u6210\u529f\u7387\u5177\u6709\u53c2\u8003\u4ef7\u503c"),
    ],
    GREEN, insert_idx=18)

prs.save("docs/\u7b54\u8fa9PPT.pptx")
print(f"PPT expanded: {len(prs.slides)} slides total")
