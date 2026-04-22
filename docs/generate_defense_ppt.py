"""
毕业答辩PPT自动生成脚本（2分钟视频版 · 10页精简）
课题：基于微服务架构的高校论坛系统设计与实现
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
PRIMARY_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT = RGBColor(0x34, 0xA8, 0x53)
DARK = RGBColor(0x20, 0x2C, 0x3C)
GRAY = RGBColor(0x5F, 0x6B, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
BORDER = RGBColor(0xDE, 0xE2, 0xE6)
ORANGE = RGBColor(0xFB, 0x8C, 0x00)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)

TOTAL_SLIDES = 14


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, bold=False, color=DARK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2),
                  font_name="微软雅黑", level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    p.level = level
    return p


def add_title_bar(slide, title_text, subtitle_text=None, page_num=None):
    bar = add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), PRIMARY, 0)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.35),
                    subtitle_text, font_size=14, color=RGBColor(0xBB, 0xD5, 0xFD))
    if page_num is not None:
        add_textbox(slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.5),
                    f"{page_num}/{TOTAL_SLIDES}", font_size=14, bold=True,
                    color=WHITE, alignment=PP_ALIGN.RIGHT)
    return bar


def add_transition(slide, text):
    add_shape_fill(slide, Inches(0), Inches(7.0), SLIDE_WIDTH, Inches(0.5),
                   RGBColor(0x13, 0x45, 0xB5), 0)
    add_textbox(slide, Inches(0.8), Inches(7.05), Inches(11), Inches(0.4),
                f"▸ {text}", font_size=12, bold=True,
                color=RGBColor(0xBB, 0xD5, 0xFD))


def add_card(slide, left, top, width, height, title, items, icon_color=PRIMARY):
    card = add_shape_fill(slide, left, top, width, height, WHITE, 0.05)
    card.shadow.inherit = False

    accent_bar = add_shape_fill(slide, left, top, Inches(0.06), height, icon_color, 0)

    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
                title, font_size=16, bold=True, color=icon_color)

    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55),
                                     width - Inches(0.5), height - Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.font.name = "微软雅黑"
        p.space_before = Pt(3)
        p.space_after = Pt(2)
    return card


def make_table(slide, left, top, width, row_height, headers, rows):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                         Emu(int(row_height * n_rows)))
    table = table_shape.table

    col_widths = [width // n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_BG
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.8), PRIMARY, 0)
    add_shape_fill(slide, Inches(0), Inches(3.6), SLIDE_WIDTH, Inches(0.4),
                   RGBColor(0x13, 0x45, 0xB5), 0)

    add_textbox(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.5),
                "毕业设计答辩", font_size=18, color=RGBColor(0xBB, 0xD5, 0xFD),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                "基于微服务架构的高校论坛系统\n设计与实现", font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(2.6), Inches(10), Inches(0.5),
                "—— PC 端 + 移动端双端协同", font_size=18, color=RGBColor(0xBB, 0xD5, 0xFD),
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.4),
                "本课题旨在采用 Spring Cloud 微服务 + Vue 3 双端技术，构建面向高校的论坛系统，"
                "解决传统单体论坛架构耦合、缺乏移动适配和测试不可复现的问题",
                font_size=13, color=RGBColor(0x99, 0xBB, 0xEE), alignment=PP_ALIGN.CENTER)

    info_items = [
        "答辩人：XXX　　学号：XXXXXXXXXX",
        "指导教师：XXX 副教授",
        "专业：软件工程　　学院：计算机与信息学院",
    ]
    y = Inches(4.3)
    for item in info_items:
        add_textbox(slide, Inches(3), y, Inches(7), Inches(0.4),
                    item, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
        y += Inches(0.45)

    add_textbox(slide, Inches(3), Inches(6.3), Inches(7), Inches(0.4),
                "二〇二六年五月", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


def slide_outline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "目 录", "CONTENTS")

    sections = [
        ("01", "研究背景与现状", "痛点 · 现有不足 · 研究现状 · 切入点"),
        ("02", "需求分析", "8个模块 · 5维非功能 · 角色权限"),
        ("03", "系统设计与架构", "四层架构 · 7子模块 · 数据库 · 选型"),
        ("04", "实现与关键技术", "认证网关 · 论坛缓存 · 双端前端"),
        ("05", "系统测试", "11场景 · 22接口 · 缓存基准 · 并发负载"),
        ("06", "总结与展望", "核心成果 · 存在问题 · 未来方向"),
    ]

    start_x = Inches(0.8)
    start_y = Inches(1.6)
    card_w = Inches(3.8)
    card_h = Inches(1.6)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)

    for i, (num, title, desc) in enumerate(sections):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_shape_fill(slide, x, y, card_w, card_h, PRIMARY_LIGHT, 0.05)

        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25),
                                            Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY
        num_circle.line.fill.background()
        tf = num_circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = "微软雅黑"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(1.0), y + Inches(0.3), Inches(2.5), Inches(0.4),
                    title, font_size=18, bold=True, color=PRIMARY)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), card_w - Inches(0.4), Inches(0.5),
                    desc, font_size=12, color=GRAY)


def slide_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "01 研究背景与意义", "Research Background & Significance")

    add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.6),
             "📋 行业与社会痛点", [
                 "• 即时通讯中课程答疑几小时后即被新消息覆盖",
                 "• 高校论坛多为单体架构，并发瓶颈、扩展困难",
                 "• 多数仅PC端，移动端用户占比已达99.6%(CNNIC)",
                 "• 功能验证依赖手工操作，交付质量缺少量化判据",
             ], PRIMARY)

    add_card(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(2.6),
             "⚠️ 现有方案的不足", [
                 "• 单体架构 → 前后端耦合、局部修改需全局重部署",
                 "• 仅PC端 → 移动端体验缺失(如水木BBS、未名BBS)",
                 "• 权限管理与内容审核机制不完善",
                 "• 缺乏自动化回归测试和性能测试",
                 "• 源码冗余度高，维护成本持续增加",
             ], ORANGE)

    add_card(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(2.6),
             "💡 课题创新点与价值", [
                 "① 微服务拆分（5个业务服务+网关） → 松耦合、可独立部署与扩展",
                 "② PC + 移动双端协同（Vue 3 + Element Plus / Vant） → 统一API、一致体验",
                 "③ JWT + Redis会话 + 网关统一鉴权 → 完整安全闭环（登录/刷新/登出/锁定）",
                 "④ Redis缓存层 → 帖子列表缓存命中后响应提升 80%+",
                 "⑤ 可复现测试链路 → 功能测试覆盖11场景22接口，通过率100%",
             ], ACCENT)


def slide_research_status(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "01 研究现状", "Related Work & Research Status")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.5),
             "📜 论坛系统技术演进", [
                 "早期：JSP/PHP 单体 → 前后端耦合、全局重部署",
                 "中期：前后端分离 → Spring Boot + Vue，效率改善[5]",
                 "当前：微服务架构 → 服务注册、网关路由、容器部署",
                 "趋势：Docker + K8s + Redis缓存 + MQ消息队列[6][21]",
             ], PRIMARY)

    add_card(slide, Inches(6.7), Inches(1.4), Inches(6.0), Inches(2.5),
             "📚 相关研究与本文定位", [
                 "微服务治理: Spring Cloud体系已形成完整技术栈[6]",
                 "安全认证: Aldea等指出JWT短生命周期+最小权限",
                 "  是缩小攻击面的关键手段[7]",
                 "迁移策略: Agaev指出单体→微服务需分阶段重构[20]",
                 "高校实践: 协同人员管理[8]、科研管理[9]等案例",
             ], ACCENT)

    add_card(slide, Inches(0.4), Inches(4.2), Inches(12.3), Inches(2.7),
             "🎯 本文切入点", [
                 "现状：论坛技术选型已从单体演进到前后端分离+微服务阶段",
                 "不足：高校场景中将微服务架构完整应用于论坛系统，并配套可量化、可复现的测试与运维体系的工程实践仍然较少",
                 "本文：在同一工程中完成从微服务拆分到双端前端适配、从认证闭环到自动化测试的完整链路",
                 "区别于：多数现有工作停留在单体框架或仅完成前后端分离",
             ], ORANGE)


def slide_requirements(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "02 需求分析与目标", "Requirements Analysis & Objectives")

    add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(3.0),
             "🔧 功能性需求（8个模块）", [
                 "用户认证：注册/登录/令牌刷新/登出/失败锁定",
                 "帖子管理：分区浏览/发帖/编辑/删除/详情",
                 "评论互动：一级评论/二级回复/编辑/删除",
                 "点赞：去重点赞(联合唯一约束)",
                 "通知服务：系统公告/个人通知/标记已读",
                 "搜索/媒体/网关访问控制",
             ], PRIMARY)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(3.0),
             "📊 非功能性需求（5个维度）", [
                 "性能：缓存命中后毫秒级响应, 未命中≤500ms",
                 "安全：JWT双令牌(Access 24h/Refresh 7d)+失败锁定",
                 "可扩展：Nacos注册发现, 新模块即插即用",
                 "可维护：Docker Compose编排+一键启停+自动回归",
                 "兼容：PC(Element Plus)+移动(Vant), 同API网关",
             ], ACCENT)

    add_card(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3),
             "🎯 论文六项研究内容", [
                 "①后端微服务拆分(6服务+Nacos+Gateway)  ②JWT+Redis认证闭环  ③论坛核心数据建模+Redis缓存",
                 "④Vue3双端前端(Element Plus/Vant)  ⑤通知/搜索/媒体辅助服务协同  ⑥Docker+脚本+Python自动化测试体系",
             ], PRIMARY)


def slide_role_permission(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "02 需求分析 — 用户角色与权限", "User Roles & Permissions")

    headers = ["功能模块", "操作", "游客", "注册用户"]
    rows = [
        ["帖子管理", "浏览列表/详情", "✅", "✅"],
        ["帖子管理", "发布帖子", "❌→登录", "✅"],
        ["帖子管理", "编辑/删除帖子", "❌", "仅作者"],
        ["评论互动", "浏览评论", "✅", "✅"],
        ["评论互动", "发表/回复评论", "❌→登录", "✅"],
        ["评论互动", "编辑/删除评论", "❌", "仅作者"],
        ["点赞", "对帖子点赞", "❌→登录", "✅(去重)"],
        ["通知服务", "首页通知摘要", "✅", "✅"],
        ["通知服务", "个人通知/标记已读", "❌→登录", "✅"],
        ["搜索/媒体", "搜索帖子/浏览推荐", "✅", "✅"],
        ["个人中心", "查看账户信息", "❌→登录", "✅"],
    ]

    make_table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Emu(int(Inches(0.40))),
               headers, rows)

    add_card(slide, Inches(0.4), Inches(5.8), Inches(12.3), Inches(1.2),
             "💡 权限模型说明", [
                 "系统不区分学生/教师角色，所有注册用户享有相同操作权限（简化的非RBAC模型）",
                 "编辑/删除权限取决于当前用户是否为内容作者 · 分区和公告通过种子数据初始化",
             ], PRIMARY)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "03 系统设计与架构 — 整体架构", "System Architecture Overview")

    layers = [
        ("用户接入层", Inches(1.3), PRIMARY_LIGHT, PRIMARY,
         ["PC端 :5173 (Vue3+Element Plus)", "移动端 :5174 (Vue3+Vant)"], Inches(3.5)),
        ("网关层", Inches(2.55), RGBColor(0xFE, 0xF3, 0xE2), ORANGE,
         ["Spring Cloud Gateway :8080 — 统一路由 · JWT验证 · CORS · 会话校验"], Inches(8.5)),
        ("业务服务层", Inches(3.8), RGBColor(0xE8, 0xF5, 0xE9), ACCENT,
         ["auth :8081", "forum :8082", "notification :8083", "media :8084", "search :8085"], Inches(8.5)),
        ("注册/配置", Inches(5.05), RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x9C, 0x27, 0xB0),
         ["Nacos 2.3.2 :8848 — 服务注册与发现 · 配置管理 · 支持直连模式降级"], Inches(8.5)),
        ("基础设施层", Inches(6.3), RGBColor(0xE3, 0xF2, 0xFD), PRIMARY,
         ["MySQL 8.4 :3306（持久化）", "Redis 7.4 :6379（缓存+会话）", "Docker Compose编排"], Inches(8.5)),
    ]

    for name, y, bg_color, border_color, items, width in layers:
        x = Inches(2.2)
        h = Inches(1.05)
        card = add_shape_fill(slide, x, y, width, h, bg_color, 0.03)
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)

        add_textbox(slide, Inches(0.5), y + Inches(0.2), Inches(1.5), Inches(0.5),
                    name, font_size=13, bold=True, color=border_color, alignment=PP_ALIGN.RIGHT)

        items_text = "　｜　".join(items)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.25), width - Inches(0.6), Inches(0.6),
                    items_text, font_size=13, color=DARK, alignment=PP_ALIGN.CENTER)

    for y_top in [Inches(2.35), Inches(3.55), Inches(4.8), Inches(6.05)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), y_top, Inches(0.35), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()


def slide_architecture_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "03 系统设计与架构 — 核心模块设计", "Core Module Design")

    modules = [
        ("common-lib", "公共库", [
            "不独立运行",
            "ApiResponse统一封装",
            "ServiceConstants全局常量",
            "JWT/Redis键前缀/请求头名",
        ], GRAY),
        ("gateway-service", "API网关", [
            "端口: 8080",
            "JWT验证+Redis会话校验",
            "全局CORS + 路由(lb://)",
            "JwtRelayFilter身份透传",
        ], PRIMARY),
        ("auth-user-service", "认证用户", [
            "端口: 8081 · 6个接口",
            "注册/登录/刷新/登出/me",
            "双令牌(Access+Refresh)",
            "失败锁定(5次/600秒)",
        ], ACCENT),
        ("forum-service", "论坛核心", [
            "端口: 8082 · 10个接口",
            "帖子/评论/点赞CRUD",
            "分区管理+种子数据",
            "Redis缓存(@Cacheable)",
        ], ORANGE),
        ("notification-service", "通知服务", [
            "端口: 8083 · 3个接口",
            "ANNOUNCEMENT/USER类型",
            "已读标记",
            "NotificationDataInitializer",
        ], RGBColor(0x9C, 0x27, 0xB0)),
        ("search-service", "搜索服务", [
            "端口: 8085 · 1个接口",
            "OpenFeign→forum-service",
            "内存过滤(title/content/",
            "  authorName/sectionName)",
        ], RED_ACCENT),
    ]

    start_x = Inches(0.4)
    start_y = Inches(1.4)
    card_w = Inches(4.0)
    card_h = Inches(2.7)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for i, (name, title, items, color) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_shape_fill(slide, x, y, card_w, card_h, WHITE, 0.04)
        card.line.color.rgb = BORDER
        card.line.width = Pt(1)

        header_bar = add_shape_fill(slide, x, y, card_w, Inches(0.55), color, 0.04)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.08), card_w - Inches(0.3), Inches(0.4),
                    f"{title} ({name})", font_size=12, bold=True, color=WHITE)

        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.2), y + Inches(0.65 + j * 0.45),
                        card_w - Inches(0.4), Inches(0.4),
                        f"▸ {item}", font_size=11, color=DARK)

    add_textbox(slide, Inches(0.4), Inches(6.6), Inches(12), Inches(0.4),
                "※ media-service (端口8084, 2个接口): 音乐/书籍推荐, 静态内存数据, 不访问MySQL/Redis — 系统中最轻量的模块",
                font_size=11, color=GRAY)


def slide_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "03 系统设计与架构 — 数据库设计", "Database Design")

    tables_data = [
        ("user_account", "用户表", [
            ("id", "BIGINT", "主键, 自增"),
            ("username", "VARCHAR(64)", "非空, 唯一"),
            ("password", "VARCHAR(255)", "BCrypt加密"),
            ("nickname", "VARCHAR(64)", "非空"),
            ("created_at", "DATETIME", "注册时间"),
            ("updated_at", "DATETIME", "更新时间"),
        ]),
        ("forum_post", "帖子表", [
            ("id", "BIGINT", "主键, 自增"),
            ("title", "VARCHAR(128)", "帖子标题"),
            ("content", "TEXT", "帖子正文"),
            ("section_id", "BIGINT", "所属分区ID"),
            ("author_id/name", "BIGINT/VARCHAR(64)", "作者(冗余)"),
            ("like_count", "INT", "默认0, 点赞计数"),
        ]),
        ("post_comment", "评论表", [
            ("id", "BIGINT", "主键, 自增"),
            ("content", "TEXT", "评论内容"),
            ("post_id", "BIGINT", "所属帖子ID"),
            ("parent_comment_id", "BIGINT", "可空(二级回复)"),
            ("user_id/username", "BIGINT/VARCHAR(64)", "评论者"),
            ("created/updated", "DATETIME", "时间戳"),
        ]),
    ]

    x_positions = [Inches(0.3), Inches(4.5), Inches(8.7)]
    y = Inches(1.4)
    card_w = Inches(4.0)

    for idx, (table_name, display_name, fields) in enumerate(tables_data):
        x = x_positions[idx]

        header = add_shape_fill(slide, x, y, card_w, Inches(0.5), PRIMARY, 0.03)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.4),
                    f"📦 {display_name} ({table_name})", font_size=13, bold=True, color=WHITE)

        for fi, (fname, ftype, fdesc) in enumerate(fields):
            fy = y + Inches(0.55 + fi * 0.38)
            row_bg = WHITE if fi % 2 == 0 else LIGHT_BG
            row = add_shape_fill(slide, x, fy, card_w, Inches(0.36), row_bg, 0)
            row.line.fill.background()
            add_textbox(slide, x + Inches(0.1), fy + Inches(0.02), Inches(1.4), Inches(0.3),
                        fname, font_size=10, bold=True, color=PRIMARY)
            add_textbox(slide, x + Inches(1.5), fy + Inches(0.02), Inches(1.1), Inches(0.3),
                        ftype, font_size=9, color=GRAY)
            add_textbox(slide, x + Inches(2.6), fy + Inches(0.02), Inches(1.3), Inches(0.3),
                        fdesc, font_size=9, color=DARK)

    other_y = Inches(4.4)
    add_card(slide, Inches(0.3), other_y, Inches(3.8), Inches(2.5),
             "📦 其他3张表 + 设计说明", [
                 "forum_section — 分区(种子数据初始化)",
                 "post_like — 点赞(联合唯一约束 post+user)",
                 "notification_message — 通知(ANNOUNCEMENT/USER)",
                 "",
                 "⚙ 6张表满足第三范式(3NF)",
                 "⚙ forum_post冗余author_name/section_name",
                 "  → 列表查询免JOIN, 配合Redis缓存降压",
             ], PRIMARY)

    add_card(slide, Inches(4.5), other_y, Inches(4.0), Inches(2.5),
             "🔗 Redis 数据结构", [
                 "auth:session:{sessionId}",
                 "  → 用户会话信息(JWT签发校验)",
                 "auth:login:fail:{username}:{ip}",
                 "  → 登录失败计数(防暴力破解)",
                 "forum:post:list::*",
                 "  → 帖子列表缓存(TTL=2min)",
                 "forum:post:detail::*",
                 "  → 帖子详情缓存(TTL=5min)",
             ], RED_ACCENT)

    add_card(slide, Inches(8.8), other_y, Inches(4.0), Inches(2.5),
             "🔒 安全机制", [
                 "密码存储: BCrypt 单向散列",
                 "令牌: JWT Access + Refresh 双令牌",
                 "会话: Redis 集中管理, 登出即失效",
                 "网关: 统一 JWT 校验 + 路径白名单",
                 "锁定: 登录失败N次后短时锁定",
                 "权限: 帖子仅作者可编辑/删除",
             ], ACCENT)


def slide_tech_selection(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "03 系统设计与架构 — 关键技术选型", "Technology Selection")

    headers = ["技术领域", "本项目选型", "对比方案", "选型理由"]
    rows = [
        ["后端框架", "Spring Boot 3.2", "Django / Express", "Java生态成熟, Spring Cloud微服务生态完善"],
        ["微服务治理", "Spring Cloud Alibaba", "Dubbo / gRPC", "Nacos集成注册+配置, 与Gateway无缝衔接"],
        ["API网关", "Spring Cloud Gateway", "Nginx / Kong", "原生支持服务发现, 与JWT Filter深度集成"],
        ["前端框架", "Vue 3 + TypeScript", "React / Angular", "组合式API, 生态轻量, 双端复用路由设计"],
        ["PC端UI", "Element Plus", "Ant Design Vue", "组件丰富, 文档完善, 社区活跃"],
        ["移动端UI", "Vant 4", "uni-app / Ionic", "Vue3原生适配, 轻量高性能"],
        ["数据库", "MySQL 8.4", "PostgreSQL / MongoDB", "事务支持完善, 运维成本低, JPA兼容性好"],  # pragma: allowlist secret
        ["缓存", "Redis 7.4", "Memcached / Caffeine", "支持多数据结构, 兼顾会话管理与数据缓存"],
        ["认证", "JWT + Spring Security", "OAuth2 / Session", "无状态令牌+Redis会话=安全与性能兼顾"],
        ["服务注册", "Nacos 2.3", "Eureka / Consul", "支持注册发现+配置管理, 阿里生态维护活跃"],
    ]

    make_table(slide, Inches(0.4), Inches(1.4), Inches(12.5), Emu(int(Inches(0.42))),
               headers, rows)


def slide_implementation_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "04 实现与关键技术 — 认证与网关", "Authentication & Gateway Implementation")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5),
             "🔐 JWT 双令牌认证流程", [
                 "1. 用户提交 username + password",
                 "2. auth-service 验证 → 签发 Access Token (短期)",
                 "                     + Refresh Token (长期)",
                 "3. Redis 写入 session: auth:session:{sid}",
                 "4. 客户端携带 Bearer Token 访问受保护接口",
                 "5. Gateway JwtRelayFilter 校验:",
                 "   · JWT签名有效性 + 类型(access)",
                 "   · Redis session存在性",
                 "   · 提取 userId/username → X-User-Id/X-Username",
                 "6. 下游服务通过请求头获取用户身份",
                 "",
                 "Refresh 流程:",
                 "   · Access过期 → 前端401拦截器自动调用 /refresh",
                 "   · 验证Refresh Token → 签发新Access Token",
                 "   · 并发请求通过 pending queue 避免重复刷新",
                 "",
                 "Logout 流程:",
                 "   · 调用 /logout → Redis 删除 session",
                 "   · 网关拦截后续请求(session不存在) → 401",
             ], PRIMARY)

    add_card(slide, Inches(6.9), Inches(1.4), Inches(6.0), Inches(2.5),
             "🌐 网关路由配置", [
                 "/api/auth/**   → auth-user-service (StripPrefix=1)",
                 "/api/users/**  → auth-user-service (StripPrefix=1)",
                 "/api/forum/**  → forum-service (StripPrefix=2)",
                 "/api/notifications/** → notification-service",
                 "/api/media/**  → media-service (StripPrefix=1)",
                 "/api/search/** → search-service (StripPrefix=1)",
             ], ACCENT)

    add_card(slide, Inches(6.9), Inches(4.2), Inches(6.0), Inches(2.7),
             "🛡️ 安全策略伪代码", [
                 "function JwtRelayFilter(request):",
                 "  if path in PUBLIC_PATHS: forward()",
                 "  token = extractBearer(request)",
                 "  claims = jwt.verify(token, SECRET)",
                 "  assert claims.type == 'access'",
                 "  session = redis.get('auth:session:' + claims.sid)",
                 "  assert session != null",
                 "  request.addHeader('X-User-Id', claims.uid)",
                 "  forward(request)",
             ], RED_ACCENT)


def slide_implementation_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "04 实现与关键技术 — 论坛核心与缓存", "Forum Core & Cache Implementation")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(3.0),
             "📝 论坛核心功能伪代码", [
                 "// 发帖流程",
                 "function createPost(request, user):",
                 "  validate(request.title, request.content)",
                 "  section = sectionRepo.findById(request.sectionId)",
                 "  post = new ForumPost(title, content, section, user)",
                 "  postRepo.save(post)",
                 "  cacheEvict('forum:post:list')  // 清除列表缓存",
                 "  return post",
                 "",
                 "// 二级评论",
                 "function addComment(postId, request, user):",
                 "  post = postRepo.findById(postId)",
                 "  comment = new PostComment(content, post, user)",
                 "  if request.parentId: comment.parent = parentComment",
                 "  commentRepo.save(comment)",
             ], PRIMARY)

    add_card(slide, Inches(6.9), Inches(1.4), Inches(6.0), Inches(3.0),
             "⚡ Redis 缓存策略", [
                 "缓存机制: Spring Cache + RedisCacheManager",
                 "",
                 "@Cacheable('forum:post:list')",
                 "function listPosts(page, size, sectionId):",
                 "  // 首次请求: 查DB → 写Redis (TTL=2min)",
                 "  // 后续请求: 直接从Redis读取",
                 "  return postRepo.findAll(pageable)",
                 "",
                 "@CacheEvict('forum:post:list', allEntries=true)",
                 "function createPost(...)  // 写操作清除缓存",
                 "",
                 "序列化: GenericJackson2Json (支持泛型反序列化)",
                 "TTL 配置: list=2min, detail=5min, default=10min",
             ], ACCENT)

    add_card(slide, Inches(0.4), Inches(4.7), Inches(6.2), Inches(2.2),
             "💡 难点与解决方案", [
                 "① 并发点赞去重 → 数据库唯一约束(post_id+user_id)",
                 "   + 业务层异常捕获返回友好提示",
                 "② 双端认证一致性 → 前端Axios 401拦截器 + pending queue",
                 "   单次刷新 + 并发请求排队重试",
                 "③ 缓存一致性 → 写操作主动 @CacheEvict + 短TTL兜底",
                 "④ 服务间调用 → OpenFeign声明式客户端(search→forum)",
             ], ORANGE)

    add_card(slide, Inches(6.9), Inches(4.7), Inches(6.0), Inches(2.2),
             "🔧 开发环境（论文4.5节）", [
                 "后端: Java 17 + Spring Boot 3.2.5 + Maven",
                 "微服务: Spring Cloud 2023.0.1 + Alibaba 2023.0.1.2",
                 "前端: Vue 3.5 + Vite 5.4 + TypeScript 5.9",
                 "UI: Element Plus 2.13 / Vant 4.9",
                 "容器: Docker Compose (MySQL8.4+Redis7.4+Nacos2.3.2)",
                 "测试: Python 3 标准库(无第三方依赖)",
             ], RGBColor(0x9C, 0x27, 0xB0))


def slide_frontend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "04 实现与关键技术 — 双端前端实现", "Dual-End Frontend Implementation")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5),
             "🖥️ PC端 (Vue3 + Element Plus)", [
                 "路由: 6个页面 (登录/帖子列表/详情/发帖/个人/通知)",
                 "",
                 "帖子列表页:",
                 "  · 分区筛选下拉 + 关键词搜索",
                 "  · 分页控制 (5/10/20条切换)",
                 "  · 侧边栏: 通知摘要 + 音乐/书籍推荐",
                 "",
                 "帖子详情页:",
                 "  · 帖子内容 + 点赞按钮",
                 "  · 评论列表 (支持二级回复)",
                 "",
                 "认证流程:",
                 "  · 顶部导航栏显示用户名/登录按钮",
                 "  · localStorage 持久化 (forum_pc_*)",
                 "  · 401自动刷新 + 登出状态同步",
                 "",
                 "状态管理: Pinia (仅auth store, 其余为组件状态)",
             ], PRIMARY)

    add_card(slide, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.5),
             "📱 移动端 (Vue3 + Vant 4)", [
                 "路由: 相同6个页面, 适配移动端布局",
                 "",
                 "交互差异:",
                 "  · NavBar导航替代顶部菜单",
                 "  · Cell组件展示帖子列表",
                 "  · 底部区域评论输入",
                 "  · 个人页集成登出 + 通知入口",
                 "",
                 "统一API层:",
                 "  · 相同的 api/ 模块 (auth/forum/notification/media)",
                 "  · 相同的 axios 实例 + 拦截器逻辑",
                 "  · 不同的 localStorage 前缀 (forum_mobile_*)",
                 "    → 双端可同域运行不冲突",
                 "",
                 "共通设计:",
                 "  · 加载态 / 空状态 / 错误重试",
                 "  · 公共浏览无需登录 (列表/详情)",
                 "  · 写操作路由守卫 (无token → /login)",
             ], ACCENT)


def slide_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "05 系统测试", "System Testing")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(3.0),
             "🧪 功能测试 (11场景 · 22接口)", [
                 "✅ 1. 获取分区列表",
                 "✅ 2. 注册并登录两个用户",
                 "✅ 3. 发帖→评论→搜索 全链路",
                 "✅ 4. 通知首页与媒体首页",
                 "✅ 5. 个人通知与标记已读",
                 "✅ 6. 无Token写操作拦截 (code=4010)",
                 "✅ 7. 无权限操作拦截 (code=4031)",
                 "✅ 8. 并发重复点赞 ([0,4091])",
                 "✅ 9. 登出后会话立即失效",
                 "✅ 10. 令牌刷新流程",
                 "✅ 11. 令牌过期拦截",
                 "全部22个API接口覆盖, 通过率100%",
             ], ACCENT)

    headers = ["测试指标", "测试结果", "说明"]
    rows = [
        ["帖子列表冷请求", "~22-31 ms", "首次请求(无缓存)"],
        ["帖子列表热请求均值", "~4-5 ms", "后续请求(Redis缓存)"],
        ["热请求 P95", "~5-6 ms", "95%请求在此延迟内"],
        ["缓存提升比例", "81-85%", "相对冷请求提升"],
        ["功能测试通过率", "100%", "11场景·22接口"],
        ["并发负载测试", "10/50/100用户", "验证承载能力"],
    ]

    make_table(slide, Inches(6.5), Inches(1.4), Inches(6.3), Emu(int(Inches(0.42))),
               headers, rows)

    add_card(slide, Inches(0.4), Inches(4.7), Inches(5.8), Inches(2.2),
             "🖥️ 测试环境", [
                 "CPU: AMD Ryzen 7 5800H · RAM: 16GB",
                 "OS: Windows 11 (64位) · JDK 17.0.8",
                 "浏览器: Microsoft Edge 146.0",
                 "脚本: Python 3 标准库(无第三方依赖)",
             ], RGBColor(0x9C, 0x27, 0xB0))

    headers2 = ["模块", "接口数", "通过数", "通过率"]
    rows2 = [
        ["认证(auth)", "6", "6", "100%"],
        ["论坛(forum)", "10", "10", "100%"],
        ["通知", "3", "3", "100%"],
        ["搜索", "1", "1", "100%"],
        ["媒体", "2", "2", "100%"],
        ["合计", "22", "22", "100%"],
    ]
    make_table(slide, Inches(6.5), Inches(4.7), Inches(6.3), Emu(int(Inches(0.30))),
               headers2, rows2)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "06 总结与展望", "Summary & Future Work")

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(3.2),
             "🏆 核心研究成果", [
                 "① 微服务架构设计与拆分",
                 "   · 5个业务服务+网关+common-lib, 四层架构",
                 "   · JWT+Redis双重鉴权, 登录失败锁定",
                 "",
                 "② 论坛核心业务 + 双端前端",
                 "   · 6张表, @Cacheable缓存, 点赞联合唯一约束",
                 "   · Vue3双端共享API层, 401静默刷新",
                 "",
                 "③ 可复现的工程验证体系",
                 "   · 11场景22接口功能测试, 通过率100%",
                 "   · 缓存基准+并发负载+Docker一键启停",
             ], ACCENT)

    add_card(slide, Inches(6.7), Inches(1.4), Inches(6.0), Inches(3.2),
             "⚠️ 存在的问题", [
                 "① 搜索服务基于内存遍历+contains匹配",
                 "   帖子增长后效率下降, 不支持分词和相关度排序",
                 "",
                 "② 日志采集与监控机制尚未建立",
                 "   仅Spring Boot默认控制台日志, 缺乏统一收集平台",
                 "",
                 "③ 部署架构未在云环境下生产级验证",
                 "   全部测试在本地单机完成",
                 "",
                 "④ 安全测试和单元测试覆盖度有提升空间",
             ], ORANGE)

    add_card(slide, Inches(0.4), Inches(4.9), Inches(12.3), Inches(2.0),
             "🚀 未来优化方向", [
                 "① 引入 Elasticsearch 全文搜索 → 倒排索引+分词分析器, 提升检索精度和速度",
                 "② ELK日志采集栈 + 链路追踪 → 实现分布式系统可观测性",
                 "③ Kubernetes集群部署 + CDN → 弹性伸缩、滚动更新、静态资源加速",
                 "④ OWASP ZAP安全扫描 + JUnit单元测试 → 完善质量保障体系",
             ], PRIMARY)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY, 0)

    add_textbox(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.5),
                "感谢各位老师的指导与聆听", font_size=40, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.0),
                "恳请各位老师批评指正", font_size=24,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(5.0),
                                  Inches(3.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xBB, 0xD5, 0xFD)
    line.line.fill.background()

    add_textbox(slide, Inches(3), Inches(5.3), Inches(7), Inches(0.5),
                "答辩人：XXX　　指导教师：XXX", font_size=16,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.5),
                "二〇二六年五月", font_size=14,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)


def slide_bg_why(prs):
    """P3: 为什么做 — 研究背景与意义"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "一、研究背景与意义", "为什么做？— 痛点、不足与研究定位", page_num=3)

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.8),
             "📋 高校论坛现状与痛点", [
                 "• 即时通讯→信息几小时后即被覆盖，难以回溯",
                 "• 早期校园BBS（水木清华/北大未名）→ 单体架构老化",
                 "• 多数仅PC端，移动端用户比例已达 99.6%（CNNIC）",
                 "• 功能验证依赖手工操作，交付质量缺少量化判据",
                 "• 前后端耦合→局部修改需全局重新部署",
             ], PRIMARY)

    add_card(slide, Inches(6.7), Inches(1.4), Inches(6.3), Inches(2.8),
             "📚 研究现状与本文定位", [
                 "• 技术演进：JSP单体 → 前后端分离[5] → 微服务[6]",
                 "• Aldea等：JWT短生命周期+最小权限缩小攻击面[7]",
                 "• Agaev：单体→微服务需分阶段重构[20]",
                 "• 高校实践：协同人员管理[8]、科研管理[9]",
                 "• 不足：论坛+微服务+完整测试的工程闭环仍较少",
             ], ACCENT)

    add_card(slide, Inches(0.4), Inches(4.5), Inches(12.6), Inches(2.5),
             "💡 课题创新点（为什么值得做）", [
                 "① 微服务拆分（5业务服务+网关）→ 松耦合、可独立部署与扩展",
                 "② PC+移动双端协同（Vue 3 + Element Plus / Vant）→ 统一API、一致体验",
                 "③ JWT+Redis会话+网关统一鉴权 → 登录/刷新/登出/锁定完整闭环",
                 "④ Redis缓存 → 帖子列表响应时间提升80%+",
                 "⑤ Python自动化测试 → 11场景22接口100%通过 + 缓存基准 + 并发负载",
             ], ORANGE)

    add_transition(slide, "明确了研究意义，接下来介绍系统的需求分析...")


def slide_req_what(prs):
    """P4: 做什么 — 需求分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "二、系统需求分析", "做什么？— 8个功能模块 · 5维非功能需求", page_num=4)

    add_card(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(2.6),
             "🔧 功能性需求（8个模块）", [
                 "用户认证：注册 / 登录 / 令牌刷新 / 登出 / 失败锁定",
                 "帖子管理：分区浏览 / 发帖 / 编辑 / 删除 / 详情",
                 "评论互动：一级评论 / 二级回复 / 编辑 / 删除",
                 "点赞去重 · 通知（公告/个人/已读）",
                 "搜索（关键词+分区）· 媒体推荐 · 网关访问控制",
             ], PRIMARY)

    add_card(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(2.6),
             "📊 非功能性需求（5维度）", [
                 "性能：缓存命中→毫秒级，未命中≤500ms",
                 "安全：JWT双令牌(24h/7d) + 失败5次锁定600s",
                 "可扩展：Nacos注册发现，新模块即插即用",
                 "可维护：Docker Compose + 一键启停 + 自动回归",
                 "兼容：PC端(Element Plus) + 移动端(Vant)",
             ], ACCENT)

    headers = ["功能模块", "操作", "游客", "注册用户"]
    rows = [
        ["帖子管理", "浏览列表/详情", "✅", "✅"],
        ["帖子管理", "发布/编辑/删除", "❌→登录", "仅作者"],
        ["评论互动", "发表/回复/删除", "❌→登录", "仅作者"],
        ["点赞", "对帖子点赞", "❌→登录", "✅(去重)"],
        ["通知/搜索/媒体", "查看/搜索/推荐", "✅(部分)", "✅"],
    ]
    make_table(slide, Inches(0.4), Inches(4.3), Inches(12.6), Emu(int(Inches(0.37))),
               headers, rows)

    add_transition(slide, "需求明确后，进入系统的架构设计...")


def slide_design_arch(prs):
    """P5: 系统设计 — 整体架构图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "三、系统设计 — 整体架构", "怎么做？— 四层架构 · 端口规划", page_num=5)

    layers = [
        ("用户接入层", Inches(1.4), PRIMARY_LIGHT, PRIMARY,
         ["PC端 :5173 (Vue3 + Element Plus)", "移动端 :5174 (Vue3 + Vant)"], Inches(4.0)),
        ("网关层", Inches(2.6), RGBColor(0xFE, 0xF3, 0xE2), ORANGE,
         ["Spring Cloud Gateway :8080 — 统一路由 · JWT验证 · CORS · Redis会话校验"], Inches(9.5)),
        ("业务服务层", Inches(3.8), RGBColor(0xE8, 0xF5, 0xE9), ACCENT,
         ["auth :8081", "forum :8082", "notification :8083", "media :8084", "search :8085"], Inches(9.5)),
        ("基础设施层", Inches(5.0), RGBColor(0xE3, 0xF2, 0xFD), PRIMARY,
         ["MySQL 8.4 :3306", "Redis 7.4 :6379", "Nacos 2.3.2 :8848", "Docker Compose"], Inches(9.5)),
    ]

    for name, y, bg_color, border_color, items, width in layers:
        x = Inches(2.0)
        h = Inches(1.0)
        card = add_shape_fill(slide, x, y, width, h, bg_color, 0.03)
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        add_textbox(slide, Inches(0.3), y + Inches(0.2), Inches(1.6), Inches(0.5),
                    name, font_size=13, bold=True, color=border_color, alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), width - Inches(0.6), Inches(0.6),
                    "　|　".join(items), font_size=13, color=DARK, alignment=PP_ALIGN.CENTER)

    for y_top in [Inches(2.4), Inches(3.6), Inches(4.8)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), y_top, Inches(0.4), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

    add_card(slide, Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.9),
             "设计依据", [
                 "各服务通过Nacos按服务名寻址 · 实例增减对上游透明 · 网关CORS+双端共用同一组API · Docker Compose一键编排",
             ], GRAY)

    add_transition(slide, "架构确定后，来看7个服务模块的具体职责...")


def slide_design_modules(prs):
    """P6: 系统设计 — 7个模块"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "三、系统设计 — 模块划分", "7个Maven子模块 · 22个API接口", page_num=6)

    modules = [
        ("common-lib", "公共库", "不独立运行 · ApiResponse统一封装\nServiceConstants全局常量", GRAY),
        ("gateway", "网关 :8080", "路由转发(lb://服务名) · JwtRelayFilter\nJWT验证+Redis会话+身份透传", PRIMARY),
        ("auth-user", "认证 :8081", "6个接口 · 注册/登录/刷新/登出/me\n双令牌 · BCrypt · 失败锁定(5次/600s)", ACCENT),
        ("forum", "论坛 :8082", "10个接口 · 帖子/评论/点赞CRUD\n分区管理 · @Cacheable缓存 · 种子数据", ORANGE),
        ("notification", "通知 :8083", "3个接口 · ANNOUNCEMENT/USER类型\n首页摘要 · 个人通知 · 标记已读", RGBColor(0x9C, 0x27, 0xB0)),
        ("search", "搜索 :8085", "1个接口 · OpenFeign→forum\n内存过滤(title/content/author/section)", RED_ACCENT),
        ("media", "媒体 :8084", "2个接口 · 音乐/书籍推荐\n静态内存数据 · 不访问DB/Redis", RGBColor(0x00, 0x97, 0xA7)),
    ]

    col_w = Inches(3.6)
    gap = Inches(0.2)

    for i, (name, title, desc, color) in enumerate(modules):
        if i < 4:
            x = Inches(0.3) + (i % 4) * (col_w + gap)
            y = Inches(1.4)
        else:
            x = Inches(0.3) + ((i - 4) % 3) * (col_w + gap) + Inches(1.9)
            y = Inches(4.2)
        h = Inches(2.5)

        card = add_shape_fill(slide, x, y, col_w, h, WHITE, 0.04)
        card.line.color.rgb = BORDER
        card.line.width = Pt(1)
        header = add_shape_fill(slide, x, y, col_w, Inches(0.5), color, 0.04)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.06), col_w - Inches(0.2), Inches(0.4),
                    f"{title} ({name})", font_size=11, bold=True, color=WHITE)

        txBox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.6),
                                         col_w - Inches(0.3), h - Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        for line in desc.split("\n"):
            p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.font.name = "微软雅黑"

    add_transition(slide, "模块划分清楚了，下面看数据库设计...")


def slide_design_db(prs):
    """P7: 系统设计 — 数据库"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "三、系统设计 — 数据库", "6张表 · 满足3NF · Redis缓存结构", page_num=7)

    headers = ["表名", "所属服务", "核心字段", "特殊约束"]
    rows = [
        ["user_account", "auth", "id, username(VARCHAR 64), password(BCrypt), nickname", "username唯一索引"],
        ["forum_section", "forum", "id, name, description, sort_order, enabled", "name唯一 · 种子初始化"],
        ["forum_post", "forum", "id, title(128), content, author_id/name, section_id/name", "冗余字段免JOIN"],
        ["post_comment", "forum", "id, post_id, content, user_id, parent_comment_id", "parent为空=一级评论"],
        ["post_like", "forum", "id, post_id, user_id, created_at", "联合唯一(post+user)"],
        ["notification_message", "notification", "id, type, user_id, title, content, is_read", "ANNOUNCEMENT/USER"],
    ]
    make_table(slide, Inches(0.3), Inches(1.3), Inches(12.7), Emu(int(Inches(0.37))),
               headers, rows)

    add_card(slide, Inches(0.3), Inches(4.5), Inches(6.2), Inches(2.7),
             "⚙ 设计要点", [
                 "• 6张表非主键字段均完全依赖主键，满足3NF",
                 "• forum_post 冗余 author_name / section_name",
                 "  → 列表查询免JOIN + Redis缓存降压",
                 "  → 代价：昵称/分区名变更时需额外同步",
                 "• post_like 联合唯一约束 uk_post_like_post_user",
                 "  → 数据库层+应用层双重去重保障",
             ], PRIMARY)

    add_card(slide, Inches(6.7), Inches(4.5), Inches(6.3), Inches(2.7),
             "🔗 Redis 数据结构", [
                 "auth:session:{sid} → 用户会话（登出即删除）",
                 "auth:login:fail:{user}:{ip} → 失败计数（原子递增）",
                 "forum:post:list → 帖子列表缓存（TTL = 2min）",
                 "forum:post:detail → 帖子详情缓存（TTL = 5min）",
                 "默认缓存 TTL = 10min",
                 "写操作 → @CacheEvict 自动淘汰相关缓存",
             ], RED_ACCENT)

    add_transition(slide, "数据模型确定后，进入核心功能的实现...")


def slide_impl_core(prs):
    """P8: 系统实现 — 认证+论坛"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "四、系统实现 — 认证与论坛核心", "伪代码 · 缓存策略 · 难点攻克", page_num=8)

    add_card(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(5.4),
             "🔐 认证伪代码", [
                 "function login(username, password):",
                 "  if locked(user,ip): throw 锁定",
                 "  user = db.findByUsername(username)",
                 "  if !bcrypt.match(password):",
                 "    redis.incr(fail_count)",
                 "    if count>=5: lock 600s",
                 "    throw 密码错误",
                 "  clearFailCount()",
                 "  sid = UUID.random()",
                 "  access = jwt.sign({uid,sid,typ:'access'})",
                 "  refresh = jwt.sign({uid,sid,typ:'refresh'})",
                 "  redis.set('auth:session:'+sid, uid)",
                 "  return {access, refresh}",
             ], PRIMARY)

    add_card(slide, Inches(4.6), Inches(1.3), Inches(4.0), Inches(5.4),
             "📝 论坛核心伪代码", [
                 "// 发帖(含缓存淘汰)",
                 "@CacheEvict('forum:post:list')",
                 "function createPost(req, user):",
                 "  section = loadSection(req.sectionId)",
                 "  post = save(title,content,section,user)",
                 "  return post",
                 "",
                 "// 点赞(双重去重)",
                 "function likePost(postId, userId):",
                 "  if exists(postId, userId):",
                 "    throw '已点赞(4091)'",
                 "  try: insert(postId, userId)",
                 "  catch UniqueViolation:",
                 "    throw '已点赞' // DB层兜底",
             ], ACCENT)

    add_card(slide, Inches(8.8), Inches(1.3), Inches(4.2), Inches(5.4),
             "💡 难点与解决", [
                 "① 并发点赞去重",
                 "  → DB联合唯一约束+业务层捕获",
                 "",
                 "② 401并发刷新",
                 "  → isRefreshing + pendingQueue",
                 "  → 单次refresh+排队重试",
                 "",
                 "③ 缓存一致性",
                 "  → @CacheEvict主动淘汰",
                 "  → 短TTL(2/5min)兜底",
                 "",
                 "④ 服务间调用",
                 "  → OpenFeign(search→forum)",
             ], ORANGE)

    add_transition(slide, "核心逻辑完成，来看双端前端的实现...")


def slide_frontend_compact(prs):
    """P9: 系统实现 — 双端前端"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "四、系统实现 — 双端前端", "PC + 移动 · 统一API · 界面展示", page_num=9)

    add_card(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(4.0),
             "🖥️ PC端 (Element Plus)", [
                 "6个页面: 登录/列表/详情/发帖/通知/个人",
                 "帖子列表: 分区筛选+搜索+分页(5/10/20)",
                 "侧边栏: 通知摘要+音乐/书籍推荐",
                 "详情页: 正文+点赞+二级评论",
                 "顶部导航栏: 用户名/登录按钮",
                 "",
                 "图 4-11 PC端帖子列表页截图",
                 "(实际演示时替换为真实截图)",
             ], PRIMARY)

    add_card(slide, Inches(4.6), Inches(1.3), Inches(4.0), Inches(4.0),
             "📱 移动端 (Vant 4)", [
                 "相同6个页面, 纵向堆叠布局",
                 "NavBar导航 + Cell列表组件",
                 "通知/推荐内嵌列表上方",
                 "固定圆形\"我\"按钮→个人中心",
                 "",
                 "统一设计:",
                 "· 加载态/空状态/错误重试",
                 "· 不同localStorage前缀(pc/mobile)",
             ], ACCENT)

    add_card(slide, Inches(8.8), Inches(1.3), Inches(4.2), Inches(4.0),
             "🔧 开发环境", [
                 "后端: Java 17 + Spring Boot 3.2.5",
                 "微服务: Spring Cloud Alibaba",
                 "前端: Vue 3.5 + Vite 5.4 + TS 5.9",
                 "UI: Element Plus 2.13 / Vant 4.9",
                 "容器: Docker Compose",
                 "测试: Python 3 标准库",
                 "",
                 "共享: Axios+Pinia+Vue Router",
                 "鉴权: 401拦截器+静默刷新",
             ], RGBColor(0x9C, 0x27, 0xB0))

    add_card(slide, Inches(0.4), Inches(5.6), Inches(12.6), Inches(1.2),
             "🔑 双端一致性设计", [
                 "共享: api/模块定义 · Pinia auth store · Vue Router路由守卫 · Axios 401拦截器+pendingQueue静默刷新",
                 "隔离: localStorage前缀(forum_pc_* / forum_mobile_*) · UI组件库(Element Plus / Vant) · 布局方式(左右分栏 / 纵向堆叠)",
             ], PRIMARY)

    add_transition(slide, "系统实现完毕，下面展示测试结果...")


def slide_impl_devops(prs):
    """P10: 系统实现 — 辅助服务+运维"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "四、系统实现 — 辅助服务与运维", "通知/搜索/媒体 · Docker · 一键启停", page_num=10)

    add_card(slide, Inches(0.4), Inches(1.4), Inches(4.0), Inches(2.8),
             "🔔 通知服务", [
                 "NotificationDataInitializer 种子公告",
                 "GET /notifications/home → 最新公告+未读数",
                 "GET /notifications/my → 个人通知列表",
                 "POST /notifications/{id}/read → 标记已读",
                 "未携带用户ID时未读数返回0",
             ], RGBColor(0x9C, 0x27, 0xB0))

    add_card(slide, Inches(4.6), Inches(1.4), Inches(4.0), Inches(2.8),
             "🔍 搜索服务", [
                 "不持有数据库 → OpenFeign调用forum",
                 "keyword为空 → 直接透传帖子列表",
                 "keyword不为空 → 内存过滤:",
                 "  title/contentPreview(前120字)/",
                 "  authorName/sectionName → contains匹配",
             ], RED_ACCENT)

    add_card(slide, Inches(8.8), Inches(1.4), Inches(4.2), Inches(2.8),
             "🎵 媒体推荐", [
                 "最轻量模块 · 不访问MySQL/Redis",
                 "推荐数据硬编码在内存中",
                 "GET /media/home → 首页摘要",
                 "GET /media/recommendations",
                 "  → 按type(music/book/all)筛选",
             ], RGBColor(0x00, 0x97, 0xA7))

    add_card(slide, Inches(0.4), Inches(4.5), Inches(12.6), Inches(2.7),
             "🐳 Docker Compose + 运维脚本", [
                 "docker-compose.yml → MySQL 8.4 + Redis 7.4 + Nacos 2.3.2（forum-net bridge网络 · 数据卷持久化 · 健康检查）",
                 "start-backend-all.sh → 启动顺序: common-lib构建 → auth → forum → notification → media → search → gateway（最后）",
                 "stop-backend-all.sh → 反向顺序关闭（网关最先停止阻止新请求）  ·  verify-backend-all.sh → /actuator/health验活",
                 "支持三种模式: full(需Nacos) / direct(无Nacos静态路由) / auto(自动判断)",
             ], PRIMARY)

    add_transition(slide, "系统实现完毕，下面展示测试结果...")


def slide_test_compact(prs):
    """P11: 系统测试"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "五、系统测试", "功能测试 · 性能测试 · 测试环境", page_num=11)

    headers = ["#", "测试场景", "预期", "结果"]
    rows = [
        ["1-5", "分区查询/注册登录/发帖评论搜索/通知媒体/已读标记", "code=0", "✅通过"],
        ["6", "无Token写操作", "code=4010", "✅通过"],
        ["7", "越权删帖", "code=4031", "✅通过"],
        ["8", "并发重复点赞", "[0,4091]", "✅通过"],
        ["9", "登出后旧Token", "code=4010", "✅通过"],
        ["10-11", "令牌刷新/令牌过期", "code=0/4010", "✅通过"],
    ]
    make_table(slide, Inches(0.3), Inches(1.3), Inches(7.5), Emu(int(Inches(0.35))),
               headers, rows)

    headers2 = ["性能指标", "数值"]
    rows2 = [
        ["冷请求", "22~31 ms"],
        ["热请求均值", "4~5 ms"],
        ["P95", "5~6 ms"],
        ["提升", "81~85%"],
        ["并发", "10/50/100用户"],
        ["接口覆盖", "22/22 (100%)"],
    ]
    make_table(slide, Inches(8.0), Inches(1.3), Inches(5.0), Emu(int(Inches(0.35))),
               headers2, rows2)

    add_card(slide, Inches(0.3), Inches(4.4), Inches(6.0), Inches(2.3),
             "🖥️ 测试环境", [
                 "CPU: AMD Ryzen 7 5800H · RAM: 16GB",
                 "OS: Windows 11 · JDK 17.0.8 · Python 3",
                 "基础设施: Docker (MySQL+Redis+Nacos)",
                 "测试方法: 全链路黑盒 · 标准库HTTP · 无第三方依赖",
             ], RGBColor(0x9C, 0x27, 0xB0))

    add_card(slide, Inches(6.5), Inches(4.4), Inches(6.5), Inches(2.3),
             "📋 三个测试脚本", [
                 "run_system_regression.py → 11场景功能回归, Markdown报告",
                 "benchmark_forum_cache.py → 冷/热响应对比, P95+提升幅度",
                 "load_test_concurrent.py → 10/50/100并发, QPS+成功率",
                 "一键执行: bash generate-test-reports.sh",
             ], PRIMARY)

    add_transition(slide, "测试全部通过，下面做系统特点总结...")


def slide_features(prs):
    """P12: 系统特点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "五、系统特点总结", "与同类毕设项目的差异化", page_num=12)

    headers = ["维度", "传统校园论坛", "本系统"]
    rows = [
        ["架构", "单体 / 前后端未分离", "微服务(5服务+网关) · Nacos注册"],
        ["终端", "仅PC端", "PC + 移动双端 · 统一API网关"],
        ["认证", "Cookie-Session", "JWT双令牌 + Redis会话 + 失败锁定"],
        ["缓存", "无 / 本地缓存", "Redis · @Cacheable · 2/5/10min TTL"],
        ["测试", "手工测试", "Python自动化 · 11场景22接口 · 缓存基准"],
        ["部署", "手工安装", "Docker Compose · 一键启停 · 三种模式"],
        ["可扩展", "修改需全局重部署", "新模块注册Nacos + 网关加路由即可"],
    ]
    make_table(slide, Inches(0.3), Inches(1.4), Inches(12.7), Emu(int(Inches(0.42))),
               headers, rows)

    add_card(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.8),
             "📌 效果数据", [
                 "功能测试: 11个端到端场景 · 22个API接口 · 通过率100%　|　缓存提升: 冷请求22~31ms → 热请求4~5ms · 提升81~85%",
                 "并发负载: 10/50/100用户 · 稳定响应 · 成功率高　|　运维闭环: 启动→验活→回归→报告 · 全流程可重复",
             ], ACCENT)

    add_transition(slide, "最后做一个全文总结...")


def slide_summary_compact(prs):
    """P13: 总结与展望"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "六、总结与展望", "核心成果 · 存在问题 · 未来方向", page_num=13)

    add_card(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(3.2),
             "🏆 三项核心成果", [
                 "① 微服务架构设计与拆分",
                 "  5个业务服务+网关+common-lib",
                 "  JWT+Redis双重鉴权, 四层架构",
                 "",
                 "② 论坛核心+双端前端",
                 "  6张表, @Cacheable缓存, 点赞去重",
                 "  Vue3双端共享API, 401静默刷新",
                 "",
                 "③ 可复现的工程验证体系",
                 "  11场景22接口100%, 缓存提升80%+",
             ], ACCENT)

    add_card(slide, Inches(4.6), Inches(1.3), Inches(4.0), Inches(3.2),
             "⚠️ 四个不足", [
                 "① 搜索基于内存遍历+contains",
                 "  不支持分词和相关度排序",
                 "",
                 "② 日志采集与监控机制未建立",
                 "  仅控制台日志, 无统一平台",
                 "",
                 "③ 未在云环境生产级验证",
                 "  全部测试在本地单机完成",
                 "",
                 "④ 安全测试和单元测试可加强",
             ], ORANGE)

    add_card(slide, Inches(8.8), Inches(1.3), Inches(4.2), Inches(3.2),
             "🚀 未来方向", [
                 "① Elasticsearch全文搜索",
                 "  倒排索引+分词, 提升检索精度",
                 "",
                 "② ELK日志栈+链路追踪",
                 "  分布式系统可观测性",
                 "",
                 "③ Kubernetes+CDN部署",
                 "  弹性伸缩, 静态资源加速",
                 "",
                 "④ OWASP ZAP+JUnit",
                 "  完善安全与单元测试",
             ], PRIMARY)

    add_card(slide, Inches(0.4), Inches(4.8), Inches(12.6), Inches(2.0),
             "📌 一句话总结", [
                 "本文围绕架构解耦、跨平台适配、安全认证和质量保障四个层面，完成了一个基于微服务架构的高校论坛系统，"
                 "验证了微服务拆分在高校论坛场景下的工程可行性，并建立了可复现的自动化测试体系。",
             ], PRIMARY)


def slide_thanks_compact(prs):
    """P14: 致谢"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_shape_fill(slide, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY, 0)

    add_textbox(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.5),
                "感谢各位老师的指导与聆听", font_size=40, bold=True,
                color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.0),
                "恳请各位老师批评指正", font_size=24,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(5.0),
                                  Inches(3.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xBB, 0xD5, 0xFD)
    line.line.fill.background()

    add_textbox(slide, Inches(3), Inches(5.3), Inches(7), Inches(0.5),
                "答辩人：XXX　　指导教师：XXX", font_size=16,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(3), Inches(5.8), Inches(7), Inches(0.5),
                "二〇二六年五月", font_size=14,
                color=RGBColor(0xBB, 0xD5, 0xFD), alignment=PP_ALIGN.CENTER)


def slide_outline_15(prs):
    """P2: 目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "目 录", "CONTENTS", page_num=2)

    sections = [
        ("01", "研究背景与意义", "为什么做？痛点 · 现状 · 创新点"),
        ("02", "系统需求分析", "做什么？功能 · 非功能 · 角色"),
        ("03", "系统设计", "怎么做？架构 · 模块 · 数据库"),
        ("04", "系统实现", "核心实现 · 双端前端 · 运维脚本"),
        ("05", "系统测试", "功能测试 · 性能测试 · 系统特点"),
        ("06", "总结与展望", "成果 · 不足 · 未来方向"),
    ]

    start_x = Inches(0.8)
    start_y = Inches(1.6)
    card_w = Inches(3.8)
    card_h = Inches(1.6)
    gap_x = Inches(0.25)
    gap_y = Inches(0.3)

    for i, (num, title, desc) in enumerate(sections):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_shape_fill(slide, x, y, card_w, card_h, PRIMARY_LIGHT, 0.05)

        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25),
                                            Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY
        num_circle.line.fill.background()
        tf = num_circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = "微软雅黑"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_textbox(slide, x + Inches(1.0), y + Inches(0.3), Inches(2.5), Inches(0.4),
                    title, font_size=18, bold=True, color=PRIMARY)

        add_textbox(slide, x + Inches(0.2), y + Inches(1.0), card_w - Inches(0.4), Inches(0.5),
                    desc, font_size=12, color=GRAY)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_cover(prs)             # P1   封面
    slide_outline_15(prs)        # P2   目录
    slide_bg_why(prs)            # P3   研究背景与意义（为什么做）
    slide_req_what(prs)          # P4   系统需求分析（做什么）
    slide_design_arch(prs)       # P5   系统设计 — 整体架构
    slide_design_modules(prs)    # P6   系统设计 — 模块划分
    slide_design_db(prs)         # P7   系统设计 — 数据库
    slide_impl_core(prs)         # P8   系统实现 — 认证与论坛核心
    slide_frontend_compact(prs)  # P9   系统实现 — 双端前端
    slide_impl_devops(prs)       # P10  系统实现 — 辅助服务与运维
    slide_test_compact(prs)      # P11  系统测试
    slide_features(prs)          # P12  系统特点总结
    slide_summary_compact(prs)   # P13  总结与展望
    slide_thanks_compact(prs)    # P14  致谢

    output_path = "/workspace/docs/毕业答辩PPT.pptx"
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")
    print()
    print("页面结构 (14页):")
    print("  P1   封面")
    print("  P2   目录")
    print("  P3   研究背景与意义（为什么做）")
    print("  P4   系统需求分析（做什么）")
    print("  P5   系统设计 — 整体架构")
    print("  P6   系统设计 — 模块划分")
    print("  P7   系统设计 — 数据库")
    print("  P8   系统实现 — 认证与论坛核心")
    print("  P9   系统实现 — 双端前端")
    print("  P10  系统实现 — 辅助服务与运维")
    print("  P11  系统测试")
    print("  P12  系统特点总结")
    print("  P13  总结与展望")
    print("  P14  致谢")
    print()
    print("结构对照:")
    print("  封面+目录          2页")
    print("  背景意义           1页  (为什么做)")
    print("  需求分析           1页  (做什么)")
    print("  系统设计           3页  (怎么做:架构/模块/数据库)")
    print("  系统实现           3页  (核心/前端/辅助运维)")
    print("  测试+特点          2页")
    print("  总结+致谢          2页")
    print("  合计              14页")


if __name__ == "__main__":
    main()
