# -*- coding: utf-8 -*-
"""Generate defense PPT v2: section title slides + content slides with animations."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from copy import deepcopy
import uuid

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BG_MID = RGBColor(0x24, 0x4E, 0x7A)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2B, 0x7A, 0xE5)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE8, 0x4D, 0x3D)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY = RGBColor(0x88, 0x88, 0x88)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                font_name=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox


def add_para(tf, text, font_size=16, bold=False, color=TEXT_DARK,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_before = space_before
    if font_name:
        p.font.name = font_name
    return p


def add_rect(slide, left, top, width, height, fill_color, text="",
             font_size=14, font_color=TEXT_DARK, bold=False):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_appear_animation(slide, shape, delay_ms=0, duration_ms=500, order=1):
    """Add a fade-in appear animation to a shape."""
    timing = slide.element.find(qn('p:timing'))
    if timing is None:
        timing_xml = (
            '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
            '<p:childTnLst/></p:cTn></p:par></p:tnLst></p:timing>'
        )
        timing = parse_xml(timing_xml)
        slide.element.append(timing)

    child_list = timing.find('.//' + qn('p:childTnLst'))
    shape_id = shape.shape_id

    anim_xml = (
        f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:cTn id="{order + 1}" fill="hold">'
        f'<p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:par><p:cTn id="{order + 2}" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:par><p:cTn id="{order + 3}" presetID="10" presetClass="entr" '
        f'presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'<p:childTnLst>'
        f'<p:set><p:cBhvr><p:cTn id="{order + 4}" dur="1" fill="hold">'
        f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        f'</p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
        f'<p:animEffect transition="in" filter="fade">'
        f'<p:cBhvr><p:cTn id="{order + 5}" dur="{duration_ms}"/>'
        f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
        f'</p:cBhvr></p:animEffect>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
        f'</p:childTnLst></p:cTn></p:par>'
    )
    child_list.append(parse_xml(anim_xml))


def make_section_slide(number, title, subtitle=""):
    """Create a section divider slide with big number + title."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, BG_MID)

    num_box = add_textbox(s, 1.0, 1.5, 3.0, 3.0, f"0{number}",
                          font_size=96, bold=True,
                          color=RGBColor(0x3D, 0x8E, 0xE8),
                          alignment=PP_ALIGN.RIGHT)

    title_box = add_textbox(s, 4.5, 2.0, 8.0, 1.5, title,
                            font_size=40, bold=True, color=TEXT_WHITE,
                            alignment=PP_ALIGN.LEFT)

    if subtitle:
        sub_box = add_textbox(s, 4.5, 3.8, 8.0, 1.0, subtitle,
                              font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE),
                              alignment=PP_ALIGN.LEFT)

    line = add_rect(s, 4.5, 3.5, 4.0, 0.06, ACCENT)

    return s


def make_content_slide(title, items, accent_color=ACCENT):
    """Create a content slide with title bar + bullet items with stagger animation."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT_BG)

    bar = add_rect(s, 0, 0, 13.333, 0.9, accent_color,
                   title, font_size=24, font_color=TEXT_WHITE, bold=True)

    shapes_to_animate = []
    for i, (heading, desc) in enumerate(items):
        y = 1.3 + i * 1.1
        h_box = add_textbox(s, 1.0, y, 11.3, 0.4, heading,
                            font_size=20, bold=True, color=accent_color)
        d_box = add_textbox(s, 1.0, y + 0.4, 11.3, 0.5, desc,
                            font_size=15, color=TEXT_GRAY)
        shapes_to_animate.append(h_box)
        shapes_to_animate.append(d_box)

    for idx, shape in enumerate(shapes_to_animate):
        add_appear_animation(s, shape, delay_ms=idx * 300,
                             duration_ms=400, order=idx * 6 + 10)

    return s


def make_two_col_slide(title, left_title, left_items, right_title, right_items,
                       left_color=ACCENT, right_color=GREEN):
    """Two-column content slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, LIGHT_BG)

    bar = add_rect(s, 0, 0, 13.333, 0.9, BG_DARK,
                   title, font_size=24, font_color=TEXT_WHITE, bold=True)

    left_bar = add_rect(s, 0.5, 1.2, 5.8, 0.5, left_color,
                        left_title, font_size=15, font_color=TEXT_WHITE, bold=True)
    right_bar = add_rect(s, 7.0, 1.2, 5.8, 0.5, right_color,
                         right_title, font_size=15, font_color=TEXT_WHITE, bold=True)

    animated = []
    for i, (h, d) in enumerate(left_items):
        y = 1.9 + i * 0.9
        hb = add_textbox(s, 0.7, y, 5.5, 0.35, h, font_size=16, bold=True, color=left_color)
        db = add_textbox(s, 0.7, y + 0.35, 5.5, 0.45, d, font_size=13, color=TEXT_GRAY)
        animated.extend([hb, db])

    for i, (h, d) in enumerate(right_items):
        y = 1.9 + i * 0.9
        hb = add_textbox(s, 7.2, y, 5.5, 0.35, h, font_size=16, bold=True, color=right_color)
        db = add_textbox(s, 7.2, y + 0.35, 5.5, 0.45, d, font_size=13, color=TEXT_GRAY)
        animated.extend([hb, db])

    for idx, shape in enumerate(animated):
        add_appear_animation(s, shape, delay_ms=idx * 200,
                             duration_ms=350, order=idx * 6 + 10)
    return s


# ============================================================
# SLIDE 1: Cover
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, BG_DARK)

t1 = add_textbox(s1, 1.0, 1.2, 11.3, 1.5,
                 "\u57fa\u4e8e\u5fae\u670d\u52a1\u67b6\u6784\u7684\n\u9ad8\u6821\u8bba\u575b\u7cfb\u7edf\u8bbe\u8ba1\u4e0e\u5b9e\u73b0",
                 font_size=38, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

t1e = add_textbox(s1, 1.0, 3.0, 11.3, 0.6,
                  "Design and Implementation of a University Forum System\nBased on Microservice Architecture",
                  font_size=14, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

line = add_rect(s1, 4.5, 3.9, 4.3, 0.04, ACCENT)

info = add_textbox(s1, 3.0, 4.2, 7.3, 1.8, "", font_size=18, color=TEXT_WHITE,
                   alignment=PP_ALIGN.CENTER)
info.text_frame.paragraphs[0].text = "\u59d3\u540d\uff1a________\u3000\u3000\u5b66\u53f7\uff1a________"
info.text_frame.paragraphs[0].font.size = Pt(20)
info.text_frame.paragraphs[0].font.color.rgb = TEXT_WHITE
info.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
add_para(info.text_frame, "\u6307\u5bfc\u8001\u5e08\uff1a\u6c6a\u7ea2\u677e", font_size=18,
         color=TEXT_WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(12))
add_para(info.text_frame, "\u534e\u5357\u5e08\u8303\u5927\u5b66 \u00b7 \u8f6f\u4ef6\u5b66\u9662",
         font_size=18, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(12))

goal = add_textbox(s1, 1.5, 6.3, 10.3, 0.6,
                   "\u201c\u9488\u5bf9\u9ad8\u6821\u8bba\u575b\u5355\u4f53\u67b6\u6784\u6269\u5c55\u6027\u4e0d\u8db3\u3001"
                   "\u7f3a\u4e4f\u79fb\u52a8\u7aef\u9002\u914d\u548c\u81ea\u52a8\u5316\u6d4b\u8bd5\u7684\u95ee\u9898\uff0c"
                   "\u8bbe\u8ba1\u5e76\u5b9e\u73b0\u4e00\u5957\u57fa\u4e8e Spring Cloud \u7684\u5fae\u670d\u52a1\u8bba\u575b\u7cfb\u7edf\u201d",
                   font_size=13, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

add_appear_animation(s1, t1, 0, 600, 10)
add_appear_animation(s1, t1e, 400, 500, 20)
add_appear_animation(s1, info, 800, 500, 30)
add_appear_animation(s1, goal, 1200, 400, 40)


# ============================================================
# SECTION 1: Background
# ============================================================
make_section_slide(1, "\u7814\u7a76\u80cc\u666f\u4e0e\u521b\u65b0\u70b9", "\u73b0\u6709\u9ad8\u6821\u8bba\u575b\u7684\u75db\u70b9\u4e0e\u672c\u6587\u7684\u6539\u8fdb\u65b9\u5411")

make_two_col_slide(
    "\u7814\u7a76\u80cc\u666f\u4e0e\u521b\u65b0\u70b9",
    "\u73b0\u6709\u9ad8\u6821\u8bba\u575b\u7684\u4e09\u4e2a\u75db\u70b9",
    [
        ("\u2776 \u5355\u4f53\u67b6\u6784\u8026\u5408", "\u5c40\u90e8\u4fee\u6539\u9700\u5168\u5c40\u91cd\u65b0\u90e8\u7f72\uff0c\u6a21\u5757\u95f4\u65e0\u6cd5\u72ec\u7acb\u6269\u5c55"),
        ("\u2777 \u7f3a\u4e4f\u79fb\u52a8\u7aef\u9002\u914d", "CNNIC \u62a5\u544a\u663e\u793a\u624b\u673a\u4e0a\u7f51\u6bd4\u4f8b\u5df2\u8fbe 99.6%"),
        ("\u2778 \u624b\u5de5\u6d4b\u8bd5\u65e0\u6cd5\u91cf\u5316", "\u529f\u80fd\u9a8c\u8bc1\u4f9d\u8d56\u4eba\u5de5\u64cd\u4f5c\uff0c\u7f3a\u5c11\u53ef\u590d\u73b0\u7684\u5224\u5b9a\u4f9d\u636e"),
    ],
    "\u672c\u6587\u7684\u4e09\u4e2a\u521b\u65b0\u70b9",
    [
        ("\u2460 \u5fae\u670d\u52a1\u67b6\u6784\u5b8c\u6574\u843d\u5730", "Nacos + Gateway + OpenFeign \u5168\u5957\u6cbb\u7406\u7ec4\u4ef6"),
        ("\u2461 JWT + Redis \u6df7\u5408\u8ba4\u8bc1", "\u4ee4\u724c\u4e3b\u52a8\u64a4\u9500 + \u524d\u7aef\u5e76\u53d1\u961f\u5217\u9759\u9ed8\u5237\u65b0"),
        ("\u2462 \u4e09\u5957\u81ea\u52a8\u5316\u6d4b\u8bd5", "\u529f\u80fd\u56de\u5f52 + \u7f13\u5b58\u57fa\u51c6 + \u5e76\u53d1\u8d1f\u8f7d\uff0c\u81ea\u52a8\u751f\u6210\u62a5\u544a"),
    ],
)


# ============================================================
# SECTION 2: Requirements
# ============================================================
make_section_slide(2, "\u9700\u6c42\u5206\u6790\u4e0e\u76ee\u6807", "\u7528\u6237\u89d2\u8272\u3001\u529f\u80fd\u6a21\u5757\u4e0e\u975e\u529f\u80fd\u6307\u6807")

make_content_slide(
    "\u9700\u6c42\u5206\u6790\u4e0e\u6280\u672f\u8def\u7ebf",
    [
        ("\u7528\u6237\u89d2\u8272", "\u6e38\u5ba2\uff08\u6d4f\u89c8/\u641c\u7d22/\u516c\u544a\uff09\u548c\u6ce8\u518c\u7528\u6237\uff08\u53d1\u5e16/\u8bc4\u8bba/\u70b9\u8d5e/\u901a\u77e5\uff09\uff0c\u7f16\u8f91\u5220\u9664\u4ec5\u9650\u4f5c\u8005"),
        ("\u516b\u5927\u529f\u80fd\u6a21\u5757", "\u8ba4\u8bc1 \u00b7 \u5e16\u5b50 \u00b7 \u8bc4\u8bba \u00b7 \u70b9\u8d5e \u00b7 \u901a\u77e5 \u00b7 \u641c\u7d22 \u00b7 \u5a92\u4f53\u63a8\u8350 \u00b7 \u7f51\u5173\u8bbf\u95ee\u63a7\u5236"),
        ("\u975e\u529f\u80fd\u6027\u9700\u6c42\u4e94\u7ef4\u6307\u6807", "\u6027\u80fd\uff08\u7f13\u5b58\u547d\u4e2d<ms\uff09\u3001\u5b89\u5168\uff08JWT+\u9501\u5b9a\uff09\u3001\u53ef\u6269\u5c55\u3001\u53ef\u7ef4\u62a4\u3001\u517c\u5bb9"),
        ("\u6280\u672f\u8def\u7ebf", "\u9700\u6c42\u5206\u6790 \u2192 \u67b6\u6784\u8bbe\u8ba1 \u2192 \u5206\u670d\u52a1\u7f16\u7801 \u2192 \u81ea\u52a8\u5316\u6d4b\u8bd5\u9a8c\u8bc1"),
    ],
    ACCENT,
)


# ============================================================
# SECTION 3: Architecture
# ============================================================
make_section_slide(3, "\u7cfb\u7edf\u8bbe\u8ba1\u4e0e\u67b6\u6784", "\u56db\u5c42\u90e8\u7f72\u3001\u6570\u636e\u5e93 ER \u56fe\u4e0e\u6280\u672f\u9009\u578b")

# Architecture layers slide
s_arch = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s_arch, LIGHT_BG)
add_rect(s_arch, 0, 0, 13.333, 0.9, BG_DARK,
         "\u56db\u5c42\u5fae\u670d\u52a1\u90e8\u7f72\u67b6\u6784", font_size=24,
         font_color=TEXT_WHITE, bold=True)

layers_data = [
    ("\u7528\u6237\u63a5\u5165\u5c42", "PC\u7aef Vue3+Element Plus :5173  |  \u79fb\u52a8\u7aef Vue3+Vant :5174",
     LIGHT_ACCENT, ACCENT),
    ("\u7f51\u5173\u5c42", "Spring Cloud Gateway :8080 \u2014 JWT\u9274\u6743 \u00b7 \u8def\u7531\u8f6c\u53d1 \u00b7 CORS",
     RGBColor(0xD5, 0xE8, 0xD4), GREEN),
    ("\u4e1a\u52a1\u670d\u52a1\u5c42", "\u8ba4\u8bc1:8081 | \u8bba\u575b:8082 | \u901a\u77e5:8083 | \u5a92\u4f53:8084 | \u641c\u7d22:8085 \u3000\u3000\u5747\u6ce8\u518c\u81f3 Nacos",
     RGBColor(0xFF, 0xF2, 0xCC), ORANGE),
    ("\u57fa\u7840\u8bbe\u65bd\u5c42", "MySQL 8.4 :3306 | Redis 7.4 :6379 | Nacos 2.3.2 :8848 \u3000\u3000Docker Compose \u7f16\u6392",
     RGBColor(0xF0, 0xF0, 0xF0), TEXT_GRAY),
]

animated_layers = []
for i, (label, content, bg, border) in enumerate(layers_data):
    y = 1.3 + i * 1.35
    rect = add_rect(s_arch, 1.0, y, 11.3, 1.1, bg)
    rect.line.color.rgb = border
    rect.line.width = Pt(2)
    add_textbox(s_arch, 1.3, y + 0.1, 2.5, 0.35, label,
                font_size=15, bold=True, color=border)
    add_textbox(s_arch, 1.3, y + 0.5, 10.8, 0.45, content,
                font_size=14, color=TEXT_DARK)
    animated_layers.append(rect)

for idx, shape in enumerate(animated_layers):
    add_appear_animation(s_arch, shape, delay_ms=idx * 400,
                         duration_ms=500, order=idx * 6 + 10)

# DB & API slide
make_two_col_slide(
    "\u6570\u636e\u5e93\u4e0e\u63a5\u53e3\u8bbe\u8ba1",
    "6 \u5f20\u8868 \u00b7 5 \u5b9e\u4f53 \u00b7 7 \u5173\u7cfb (Chen ER)",
    [
        ("user \u2192 post / comment / like / notification", "\u4e00\u5bf9\u591a\u5173\u8054\uff0c\u70b9\u8d5e\u4e3a M:N \u901a\u8fc7 post_like \u62c6\u89e3"),
        ("section \u2192 post\uff0ccomment \u81ea\u5173\u8054", "\u5206\u533a\u5305\u542b\u5e16\u5b50\uff0c\u8bc4\u8bba\u901a\u8fc7 parentCommentId \u56de\u590d"),
        ("\u6280\u672f\u9009\u578b", "Spring Boot 3 + Spring Cloud Alibaba\uff0c\u5bf9\u6bd4 Django \u7f3a\u5c11 Nacos/Gateway \u96c6\u6210"),
    ],
    "22 \u4e2a RESTful API \u63a5\u53e3",
    [
        ("\u8ba4\u8bc1\u670d\u52a1 6 \u4e2a", "\u6ce8\u518c/\u767b\u5f55/\u5237\u65b0/\u767b\u51fa/\u4e2a\u4eba\u4fe1\u606f/\u63a2\u6d4b"),
        ("\u8bba\u575b\u670d\u52a1 10 \u4e2a", "\u5e16\u5b50CRUD + \u8bc4\u8bba\u4e09\u64cd\u4f5c + \u70b9\u8d5e + \u5206\u533a"),
        ("\u8f85\u52a9\u670d\u52a1 6 \u4e2a", "\u901a\u77e5 3 + \u641c\u7d22 1 + \u5a92\u4f53 2"),
    ],
    ACCENT, GREEN,
)


# ============================================================
# SECTION 4: Implementation
# ============================================================
make_section_slide(4, "\u5b9e\u73b0\u4e0e\u5173\u952e\u6280\u672f", "\u8ba4\u8bc1\u95ed\u73af\u3001\u7f13\u5b58\u7b56\u7565\u3001\u5e76\u53d1\u63a7\u5236\u4e0e\u5de5\u5177\u94fe")

make_content_slide(
    "\u6838\u5fc3\u5b9e\u73b0\u4e0e\u96be\u70b9",
    [
        ("JWT + Redis \u8ba4\u8bc1\u95ed\u73af", "\u767b\u5f55\u7b7e\u53d1\u53cc\u4ee4\u724c\uff0cRedis \u5b58\u50a8\u4f1a\u8bdd\uff0c\u767b\u51fa\u5220\u9664\u4f1a\u8bdd\u5373\u65f6\u5931\u6548\uff0c\u7f51\u5173\u53cc\u91cd\u6821\u9a8c"),
        ("Redis \u7f13\u5b58\u53cc\u8def\u5f84", "\u8bfb\u547d\u4e2d\u76f4\u63a5\u8fd4\u56de\uff0c\u672a\u547d\u4e2d\u67e5 MySQL \u5e76\u56de\u5199\uff1b\u5199\u64cd\u4f5c\u901a\u8fc7 @CacheEvict \u81ea\u52a8\u6dd8\u6c70"),
        ("\u5e76\u53d1\u70b9\u8d5e\u53cc\u91cd\u53bb\u91cd", "\u5e94\u7528\u5c42 exists \u5feb\u901f\u68c0\u67e5 + \u6570\u636e\u5e93\u8054\u5408\u552f\u4e00\u7ea6\u675f\u5146\u5e95\uff0c\u96f6\u9501\u5f00\u9500"),
        ("\u524d\u7aef\u9759\u9ed8\u5237\u65b0", "\u54cd\u5e94\u62e6\u622a\u5668\u5904\u7406 401\uff0c\u5e76\u53d1\u8bf7\u6c42\u52a0\u5165\u7b49\u5f85\u961f\u5217\uff0c\u907f\u514d\u591a\u6b21\u89e6\u53d1\u5237\u65b0"),
        ("\u5f00\u53d1\u5de5\u5177\u94fe", "Java 17 + Maven + Spring Boot 3.2.5 | Vue 3.5 + Vite 5.4 | Docker + Bash + Python"),
    ],
    ACCENT,
)


# ============================================================
# SECTION 5: Test
# ============================================================
make_section_slide(5, "\u6d4b\u8bd5\u7ed3\u679c", "\u529f\u80fd\u56de\u5f52\u3001\u7f13\u5b58\u57fa\u51c6\u3001\u5e76\u53d1\u8d1f\u8f7d")

# Test results - three columns
s_test = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s_test, LIGHT_BG)
add_rect(s_test, 0, 0, 13.333, 0.9, BG_DARK,
         "\u6d4b\u8bd5\u7ed3\u679c\u6982\u89c8", font_size=24,
         font_color=TEXT_WHITE, bold=True)

cols = [
    (ACCENT, "\u529f\u80fd\u6d4b\u8bd5",
     ["11 \u4e2a\u7aef\u5230\u7aef\u573a\u666f", "22 \u4e2a API \u5168\u8986\u76d6", "\u901a\u8fc7\u7387 100%", "\u542b 4 \u79cd\u5f02\u5e38\u8def\u5f84\u9a8c\u8bc1"]),
    (GREEN, "\u7f13\u5b58\u57fa\u51c6",
     ["\u51b7\u542f\u52a8 59.12 ms", "\u70ed\u7f13\u5b58 14.98 ms", "P95: 26.3 ms", "\u63d0\u5347 74.66%"]),
    (ORANGE, "\u5e76\u53d1\u8d1f\u8f7d",
     ["10\u2192 QPS 1186", "100\u2192 QPS 1392, 99.95%", "500\u2192 QPS 1557", "1000\u2192 QPS 1450, 99.99%"]),
]

anim_shapes = []
for ci, (color, title, items) in enumerate(cols):
    x = 0.5 + ci * 4.2
    bar = add_rect(s_test, x, 1.2, 3.8, 0.5, color,
                   title, font_size=15, font_color=TEXT_WHITE, bold=True)
    anim_shapes.append(bar)
    box = add_textbox(s_test, x + 0.2, 1.9, 3.4, 3.5, "", font_size=15)
    box.text_frame.paragraphs[0].text = items[0]
    box.text_frame.paragraphs[0].font.size = Pt(15)
    for item in items[1:]:
        add_para(box.text_frame, item, font_size=15, space_before=Pt(12))
    anim_shapes.append(box)

for idx, shape in enumerate(anim_shapes):
    add_appear_animation(s_test, shape, delay_ms=idx * 300,
                         duration_ms=400, order=idx * 6 + 10)


# ============================================================
# SECTION 6: Conclusion
# ============================================================
make_section_slide(6, "\u603b\u7ed3\u4e0e\u5c55\u671b", "\u6838\u5fc3\u7ed3\u8bba\u3001\u5c40\u9650\u6027\u4e0e\u672a\u6765\u65b9\u5411")

make_two_col_slide(
    "\u603b\u7ed3\u4e0e\u5c55\u671b",
    "\u4e09\u70b9\u6838\u5fc3\u7ed3\u8bba",
    [
        ("\u2714 \u5fae\u670d\u52a1\u53ef\u843d\u5730", "\u4e94\u4e2a\u670d\u52a1\u72ec\u7acb\u90e8\u7f72\uff0c\u9a8c\u8bc1\u4e86\u67b6\u6784\u5728\u4e2d\u7b49\u89c4\u6a21 Web \u5e94\u7528\u4e2d\u7684\u53ef\u884c\u6027"),
        ("\u2714 \u53cc\u7aef\u9002\u914d\u53ef\u884c", "PC + \u79fb\u52a8\u7aef\u5171\u4eab API \u5c42\uff0c\u7528\u6237\u5207\u6362\u8bbe\u5907\u65f6\u4f53\u9a8c\u4e00\u81f4"),
        ("\u2714 \u6d4b\u8bd5\u53ef\u590d\u73b0", "\u4e09\u5957\u811a\u672c\u8986\u76d6\u529f\u80fd/\u7f13\u5b58/\u5e76\u53d1\uff0c\u7ed3\u679c\u53ef\u91cf\u5316"),
    ],
    "\u4e0d\u8db3\u4e0e\u5c55\u671b",
    [
        ("\u25b3 \u641c\u7d22\u80fd\u529b\u6709\u9650", "\u5185\u5b58\u8fc7\u6ee4\u65b9\u6848\uff0c\u540e\u7eed\u53ef\u5f15\u5165 Elasticsearch"),
        ("\u25b3 \u672a\u7ecf\u4e91\u7aef\u9a8c\u8bc1", "\u672c\u5730\u5355\u673a\u6d4b\u8bd5\uff0c\u540e\u7eed\u53ef\u90e8\u7f72\u81f3 K8s \u96c6\u7fa4"),
        ("", ""),
    ],
    GREEN, RED,
)


# ============================================================
# SLIDE: Thank you
# ============================================================
s_end = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s_end, BG_DARK)

thanks = add_textbox(s_end, 1.0, 2.0, 11.3, 2.0,
                     "\u8c22\u8c22\u5404\u4f4d\u8001\u5e08\uff01",
                     font_size=48, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

sub = add_textbox(s_end, 1.0, 4.2, 11.3, 1.0,
                  "\u8bf7\u8001\u5e08\u63d0\u95ee",
                  font_size=28, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

add_appear_animation(s_end, thanks, 0, 800, 10)
add_appear_animation(s_end, sub, 600, 600, 20)

# Save
prs.save("docs/\u7b54\u8fa9PPT.pptx")
print(f"PPT saved: docs/\u7b54\u8fa9PPT.pptx")
print(f"Total slides: {len(prs.slides)}")
